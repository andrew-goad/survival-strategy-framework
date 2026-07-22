"""
====================================================================================================
TITLE:              SURVIVAL STRATEGY FRAMEWORK
SUBTITLE:           From Time-to-Event Data to Governed Intervention Evidence
AUTHOR:             Andrew R. Goad
PUBLIC REPOSITORY:  github.com/andrew-goad/survival-strategy-framework
ENGINE RELEASE:     2026.07-mainline
PYTHON:             3.11+
----------------------------------------------------------------------------------------------------
FUNCTIONAL PURPOSE
----------------------------------------------------------------------------------------------------
This module implements a reusable, portfolio-grade survival-analysis decision framework. It is
intentionally designed as more than a fitted Cox proportional hazards model. The framework connects:

    governed time-to-event inputs
        -> independent persona discovery
        -> regularized Cox proportional hazards modeling
        -> cross-validated discrimination and horizon calibration
        -> proportional-hazards diagnostics and sensitivity review
        -> model-derived risk stratification
        -> same-cohort intervention / control / stress simulation
        -> dependency-safe feature reconstruction
        -> predicted survival comparison
        -> immutable executive and technical evidence

The primary demonstration uses deterministic synthetic retention data. A secondary Rossi-dataset
smoke test remains available only to demonstrate portability to another conforming survival dataset.

----------------------------------------------------------------------------------------------------
EXECUTIVE TALK TRACKS
----------------------------------------------------------------------------------------------------
1. PREDICTING THE RUNWAY
   Binary classification asks whether an event may occur. Survival analysis adds the timing
   dimension: when is risk likely to emerge, which groups have the shortest modeled runway, and
   where might a controlled intervention be most valuable?

2. PERSONAS AND MODEL RISK ARE COMPLEMENTARY — NOT THE SAME THING
   K-Means personas describe naturally occurring lifecycle profiles. The CoxPH model independently
   estimates relative event hazard. Persona labels are never inserted into the CoxPH feature matrix.
   This separation prevents descriptive segmentation from being misrepresented as a causal or model
   risk driver.

3. APPARENT FIT IS NOT VALIDATED PERFORMANCE
   The framework reports the development-sample concordance index explicitly as apparent / in-sample
   evidence. It also performs deterministic stratified K-fold cross-validation, archives fold-level
   concordance, and creates out-of-fold predicted survival for calibration testing.

4. GOVERNED, SAME-COHORT SIMULATION
   The framework selects a baseline high-risk cohort, applies controlled scenarios only to those
   governed IDs, rebuilds dependent interactions and squared terms, and re-scores the identical
   population. A no-change control and an adverse stress scenario verify neutral and bidirectional
   model sensitivity.

5. PUSHING THE CURVE TO THE RIGHT — WITH CORRECT BOUNDARIES
   Scenario evidence includes relative mean partial-hazard movement and predicted survival
   probability uplift at governed horizons. Survival uplift is reported in percentage points.
   The output is modeled sensitivity, not causal impact, booked revenue, or realized ROI.

6. NO COLD HANDOFFS
   Each run produces a reproducible evidence package containing input validation, model coefficients,
   hazard ratios, cross-validation, calibration, proportional-hazards diagnostics, persona-quality
   measures, cohort profiles, scenario-change controls, dependency audits, reproducibility checks,
   survival visuals, an executive PowerPoint, and a technical PDF.

----------------------------------------------------------------------------------------------------
STATISTICAL AND GOVERNANCE BOUNDARIES
----------------------------------------------------------------------------------------------------
- This is a survival-analysis and scenario-sensitivity framework, not a causal-inference engine.
- Scenario results describe fitted-model sensitivity under controlled feature changes. They do not
  prove that an operational treatment will cause the modeled outcome.
- The primary demonstration uses synthetic data and contains no PII or proprietary customer data.
- Partial hazard is a relative-risk quantity. It is not an event probability.
- Kaplan-Meier curves are non-parametric observed survival estimates. Baseline-versus-scenario
  curves are CoxPH model-predicted survival estimates. The two evidence types are labeled separately.
- Apparent concordance is development-sample evidence. Cross-validated concordance is the more
  appropriate generalization estimate within this demonstration.
- Calibration results use out-of-fold predictions. They remain demonstration evidence and do not
  replace external, temporal, or institution-specific validation.
- Proportional-hazards tests are reviewed at both raw and multiplicity-adjusted significance levels.
  A raw diagnostic flag can produce PASS_WITH_REVIEW rather than being silently discarded.
- The framework does not calculate booked revenue, realized ROI, or defended LTV.
- Production use requires domain-specific event definitions, censoring rules, feature governance,
  external / temporal validation, calibration monitoring, stability monitoring, and model-risk
  approval.

----------------------------------------------------------------------------------------------------
CORE DATA CONTRACT
----------------------------------------------------------------------------------------------------
Every input dataset must provide:

1. A unique ID column.
2. A strictly positive numeric duration column.
3. A binary event column where 1 = event observed and 0 = right-censored / still active.
4. Numeric or Boolean CoxPH features with no missing or infinite values.
5. Numeric or Boolean persona-segmentation features with no missing or infinite values.
6. Engineered terms that follow these conventions when used:
      - interaction: <feature_a>_x_<feature_b>[_x_<feature_c>...]
      - square:      <feature>_sq

Scenarios may change governed base features. They may not directly change the ID, duration, event,
or engineered-term columns. Engineered terms are rebuilt automatically after base-feature changes.

----------------------------------------------------------------------------------------------------
RUN OUTPUT PACKAGE
----------------------------------------------------------------------------------------------------
Each successful run creates an immutable timestamped directory containing, among other artifacts:

    run_metadata.json
    framework_config.json
    scenario_definitions.json
    input_validation.csv
    input_snapshot.csv
    acceptance_checks.csv
    reproducibility_checks.csv

    model_summary.csv
    model_fit_metadata.csv
    model_fit_warnings.csv
    cross_validation_results.csv
    cross_validated_predictions.csv
    calibration_metrics.csv
    calibration_by_risk_group.csv
    proportional_hazards_test.csv
    ph_sensitivity_results.csv
    scaled_schoenfeld_residuals.csv

    persona_quality_metrics.csv
    persona_stability_results.csv
    persona_profiles.csv
    risk_tier_profiles.csv
    target_cohort.csv

    scenario_results.csv
    scenario_change_audit.csv
    scenario_target_scores.csv
    dependency_audit.csv
    predicted_survival_curves.csv

    charts/
        persona_kaplan_meier.png
        risk_tier_kaplan_meier.png
        baseline_vs_scenario_survival.png
        scenario_hazard_reduction.png
        calibration_at_horizons.png
        ph_assumption_diagnostics.png

    Survival_Strategy_Deck.pptx
    Technical_Model_Evidence.pdf
    artifact_manifest.json

----------------------------------------------------------------------------------------------------
PRIMARY DEPENDENCIES
----------------------------------------------------------------------------------------------------
- pandas / numpy: data engineering, validation, and controlled simulation
- scikit-learn: scaling, K-Means, persona-quality metrics, and stratified cross-validation
- lifelines: CoxPH, Kaplan-Meier, concordance, predicted survival, and PH diagnostics
- matplotlib: evidence visuals
- python-pptx: executive reporting
- reportlab: technical model-evidence PDF
====================================================================================================
"""

from __future__ import annotations

# ==================================================================================================
# STANDARD LIBRARY IMPORTS
# ==================================================================================================

import argparse
import hashlib
import importlib.metadata
import json
import platform
import sys
import traceback
import warnings
from dataclasses import asdict, dataclass, field, is_dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Literal, Mapping, Sequence

# ==================================================================================================
# THIRD-PARTY ANALYTICAL IMPORTS
# ==================================================================================================

import matplotlib

# Headless rendering is required for servers, GitHub Actions, and command-line execution.
matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from sklearn.cluster import KMeans
from sklearn.metrics import (
    adjusted_rand_score,
    calinski_harabasz_score,
    davies_bouldin_score,
    silhouette_score,
)
from sklearn.model_selection import StratifiedKFold
from sklearn.preprocessing import StandardScaler

# Survival-analysis engine.
from lifelines import CoxPHFitter, KaplanMeierFitter
from lifelines.datasets import load_rossi
from lifelines.statistics import proportional_hazard_test
from lifelines.utils import concordance_index

# Executive-reporting stack.
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Technical-reporting stack.
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image as RLImage,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


# ==================================================================================================
# ENGINE CONSTANTS AND VISUAL TOKENS
# ==================================================================================================

ENGINE_RELEASE = "2026.07-mainline"

# Report colors intentionally mirror the broader portfolio's restrained dark technical aesthetic.
COLOR_BACKGROUND = "#07131A"
COLOR_PANEL = "#0C202A"
COLOR_CYAN = "#20D7E5"
COLOR_TEAL = "#1BA7A8"
COLOR_GOLD = "#D9A441"
COLOR_WHITE = "#F4F5F6"
COLOR_LIGHT_GRAY = "#C7D0D5"
COLOR_MID_GRAY = "#71808A"
COLOR_RED = "#D95C5C"
COLOR_GREEN = "#56C271"


# ==================================================================================================
# CONFIGURATION CONTRACTS
# ==================================================================================================

@dataclass(frozen=True)
class FrameworkConfig:
    """Defines the governed analytical, validation, and reporting contract for one run.

    Teaching note:
    A reusable framework should not hide important assumptions inside functions. Centralizing the
    duration field, event field, features, model controls, cross-validation, calibration horizons,
    persona controls, scenario tolerances, and output rules makes the run reproducible and reviewable.
    """

    project_name: str
    duration_col: str
    event_col: str
    id_col: str

    # Persona and CoxPH features are deliberately separate. Persona labels do not enter the model.
    model_features: tuple[str, ...]
    segmentation_features: tuple[str, ...]

    # Optional fields summarized in persona and risk-tier profile tables.
    profile_features: tuple[str, ...] = ()

    # Human-readable time unit used in charts, narrative, and scenario evidence.
    time_unit_label: str = "months"

    # Core persona / risk / model controls.
    n_clusters: int = 3
    risk_quantile: float = 0.75
    penalizer: float = 0.10
    l1_ratio: float = 0.0
    random_state: int = 42
    standardize_model_features: bool = True

    # Model-development and validation controls.
    cross_validation_folds: int = 5
    cross_validation_shuffle: bool = True
    calibration_horizons: tuple[float, ...] = (6.0, 12.0, 18.0, 24.0)
    calibration_groups: int = 10
    calibration_review_threshold: float = 0.10

    # A horizon only slightly above the largest observed duration may reflect rounding rather than
    # meaningful extrapolation. The tolerance avoids false warnings while preserving a governed
    # warning when the requested horizon extends materially beyond the observed support.
    horizon_extrapolation_tolerance_ratio: float = 0.01

    show_km_confidence_intervals: bool = True
    run_ph_test: bool = True
    ph_test_alpha: float = 0.05
    ph_strata_max_unique: int = 10

    # Persona-quality and reproducibility controls.
    persona_silhouette_sample_size: int = 2_000
    persona_silhouette_review_threshold: float = 0.10
    persona_stability_repeats: int = 5
    verify_reproducibility: bool = True

    # Scenario-evidence controls.
    evaluation_horizon: float = 12.0
    timeline_points: int = 121
    neutral_scenario_tolerance: float = 1e-10

    # Validation thresholds are warnings unless the underlying data violate a hard contract.
    low_event_rate_warning: float = 0.05
    high_event_rate_warning: float = 0.95
    high_missingness_warning: float = 0.20

    # Output governance.
    output_root: Path = Path("outputs")
    save_input_snapshot: bool = True


@dataclass(frozen=True)
class FeatureChange:
    """Defines one bounded scenario change to a governed base feature.

    operation:
        multiply -> new = old * value
        add      -> new = old + value
        replace  -> new = value

    lower_bound / upper_bound protect the feature domain. round_digits is useful for count-like
    fields such as product adoption. The scenario audit records configured versus realized movement,
    including how often a floor or cap was reached.
    """

    feature: str
    operation: Literal["multiply", "add", "replace"]
    value: float
    lower_bound: float | None = None
    upper_bound: float | None = None
    round_digits: int | None = None


@dataclass(frozen=True)
class ScenarioDefinition:
    """Defines one controlled scenario and its expected modeled direction.

    scenario_type supports four distinct analytical purposes:
      CONTROL     -> no-change benchmark; expected modeled movement is neutral
      IMPROVEMENT -> favorable retention hypothesis; expected movement is improved
      STRESS      -> adverse sensitivity; expected movement is adverse
      TECHNICAL   -> portability / plumbing test; no business direction is asserted
    """

    name: str
    description: str
    changes: tuple[FeatureChange, ...] = field(default_factory=tuple)
    scenario_type: Literal["CONTROL", "IMPROVEMENT", "STRESS", "TECHNICAL"] = "IMPROVEMENT"
    expected_direction: Literal["NEUTRAL", "IMPROVED", "ADVERSE", "UNSPECIFIED"] = "IMPROVED"


# ==================================================================================================
# RESULT AND ARTIFACT CONTAINERS
# ==================================================================================================

@dataclass(frozen=True)
class ValidationIssue:
    severity: Literal["ERROR", "WARNING", "INFO"]
    area: str
    field: str
    message: str


@dataclass
class ValidationResult:
    issues: list[ValidationIssue] = field(default_factory=list)

    def add(self, severity: str, area: str, field_name: str, message: str) -> None:
        self.issues.append(
            ValidationIssue(
                severity=severity.upper(),  # type: ignore[arg-type]
                area=area,
                field=field_name,
                message=message,
            )
        )

    @property
    def has_errors(self) -> bool:
        return any(issue.severity == "ERROR" for issue in self.issues)

    def to_frame(self) -> pd.DataFrame:
        if not self.issues:
            return pd.DataFrame(
                [{
                    "Severity": "INFO",
                    "Area": "Input Validation",
                    "Field": "ALL",
                    "Message": "All hard input-contract checks passed.",
                }]
            )
        return pd.DataFrame(
            [
                {
                    "Severity": issue.severity,
                    "Area": issue.area,
                    "Field": issue.field,
                    "Message": issue.message,
                }
                for issue in self.issues
            ]
        )

    def raise_for_errors(self) -> None:
        if not self.has_errors:
            return
        error_text = "\n".join(
            f"[{issue.area}] {issue.field}: {issue.message}"
            for issue in self.issues
            if issue.severity == "ERROR"
        )
        raise ValueError(f"Survival input validation failed:\n{error_text}")


@dataclass
class PersonaArtifacts:
    scored_data: pd.DataFrame
    scaler: StandardScaler
    kmeans: KMeans
    profiles: pd.DataFrame
    quality_metrics: pd.DataFrame
    stability_results: pd.DataFrame
    persona_name_map: Mapping[int, str]


@dataclass
class CoxModelArtifacts:
    model: CoxPHFitter
    model_scaler: StandardScaler | None
    training_frame: pd.DataFrame
    model_summary: pd.DataFrame
    fit_metadata: pd.DataFrame
    fit_warnings: pd.DataFrame

    cross_validation_results: pd.DataFrame
    cross_validated_predictions: pd.DataFrame
    calibration_metrics: pd.DataFrame
    calibration_by_risk_group: pd.DataFrame

    ph_test_results: pd.DataFrame
    ph_residuals: pd.DataFrame
    ph_sensitivity_results: pd.DataFrame


@dataclass
class RiskArtifacts:
    scored_data: pd.DataFrame
    risk_tier_profiles: pd.DataFrame
    target_cohort: pd.DataFrame
    target_index: pd.Index
    target_threshold: float
    target_label: str


@dataclass
class ScenarioArtifacts:
    summary: pd.DataFrame
    target_scores: pd.DataFrame
    scenario_change_audit: pd.DataFrame
    dependency_audit: pd.DataFrame
    predicted_survival_curves: pd.DataFrame


@dataclass
class PlotArtifacts:
    persona_km: Path
    risk_tier_km: Path
    baseline_vs_scenario: Path
    scenario_comparison: Path
    calibration: Path
    ph_diagnostics: Path


@dataclass
class RunArtifacts:
    run_id: str
    run_directory: Path
    validation: pd.DataFrame
    personas: PersonaArtifacts
    model: CoxModelArtifacts
    risk: RiskArtifacts
    scenarios: ScenarioArtifacts
    plots: PlotArtifacts
    executive_narrative: str
    acceptance_checks: pd.DataFrame
    reproducibility_checks: pd.DataFrame
    acceptance_status: str
    power_point: Path
    technical_pdf: Path
    manifest: Mapping[str, str]

# ==================================================================================================
# GENERIC UTILITY FUNCTIONS
# ==================================================================================================

def utc_now() -> datetime:
    """Returns a timezone-aware UTC timestamp for reproducible run lineage."""
    return datetime.now(timezone.utc)


def safe_slug(value: str) -> str:
    """Converts free text into a file-system-safe identifier."""
    cleaned = "".join(ch.lower() if ch.isalnum() else "-" for ch in value.strip())
    while "--" in cleaned:
        cleaned = cleaned.replace("--", "-")
    return cleaned.strip("-") or "run"


def to_jsonable(value: Any) -> Any:
    """Recursively converts dataclasses, Paths, tuples, and NumPy values into JSON-safe values."""
    if is_dataclass(value):
        return {key: to_jsonable(item) for key, item in asdict(value).items()}
    if isinstance(value, Path):
        return str(value)
    if isinstance(value, Mapping):
        return {str(key): to_jsonable(item) for key, item in value.items()}
    if isinstance(value, (tuple, list)):
        return [to_jsonable(item) for item in value]
    if isinstance(value, (np.integer,)):
        return int(value)
    if isinstance(value, (np.floating,)):
        return float(value)
    if isinstance(value, (np.bool_,)):
        return bool(value)
    if isinstance(value, datetime):
        return value.isoformat()
    return value


def write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as stream:
        json.dump(to_jsonable(payload), stream, indent=2, sort_keys=True)


def package_versions() -> dict[str, str]:
    """Captures the execution environment without failing if a package name is unavailable."""
    packages = [
        "pandas",
        "numpy",
        "matplotlib",
        "scikit-learn",
        "lifelines",
        "python-pptx",
        "reportlab",
    ]
    versions: dict[str, str] = {"python": platform.python_version()}
    for package in packages:
        try:
            versions[package] = importlib.metadata.version(package)
        except importlib.metadata.PackageNotFoundError:
            versions[package] = "NOT_INSTALLED"
    return versions


def dataframe_fingerprint(df: pd.DataFrame, id_col: str) -> str:
    """Creates a stable SHA-256 fingerprint for the run input.

    Sorting by the governed ID prevents a harmless row-order change from producing a different
    fingerprint. The fingerprint is lineage evidence; it is not a cryptographic privacy control.
    """
    ordered = df.sort_values(id_col).reset_index(drop=True)
    hashed = pd.util.hash_pandas_object(ordered, index=False).values.tobytes()
    return hashlib.sha256(hashed).hexdigest()


def create_run_directory(config: FrameworkConfig, fingerprint: str) -> tuple[str, Path]:
    timestamp = utc_now().strftime("%Y%m%dT%H%M%SZ")
    run_id = f"{timestamp}_{fingerprint[:8]}"
    run_directory = config.output_root / safe_slug(config.project_name) / run_id
    (run_directory / "charts").mkdir(parents=True, exist_ok=False)
    return run_id, run_directory


def numeric_frame(df: pd.DataFrame, columns: Sequence[str]) -> pd.DataFrame:
    """Returns a clean floating-point feature matrix while preserving row index and column names."""
    result = df.loc[:, list(columns)].copy()
    for column in result.columns:
        if pd.api.types.is_bool_dtype(result[column]):
            result[column] = result[column].astype(int)
        result[column] = pd.to_numeric(result[column], errors="raise").astype(float)
    return result


def zscore(values: np.ndarray) -> np.ndarray:
    """Stable z-score helper used only inside deterministic synthetic-data construction."""
    std = float(np.std(values))
    if std == 0:
        return np.zeros_like(values, dtype=float)
    return (values - float(np.mean(values))) / std


def format_percent(value: float, decimals: int = 2) -> str:
    if pd.isna(value):
        return ""
    return f"{value:.{decimals}%}"


def format_float(value: float, decimals: int = 3) -> str:
    if pd.isna(value):
        return ""
    return f"{value:,.{decimals}f}"


def horizon_token(value: float) -> str:
    """Creates a stable column-name token for a numeric evaluation horizon."""
    return f"{value:g}".replace("-", "neg").replace(".", "p")


def holm_adjusted_pvalues(p_values: Sequence[float]) -> np.ndarray:
    """Returns Holm family-wise-error adjusted p-values without requiring statsmodels.

    Teaching note:
    When several feature-level PH tests are run, at least one raw p-value may fall below 0.05 by
    chance. Holm adjustment is a conservative step-down correction that controls the family-wise
    error rate while preserving a visible record of the original raw diagnostic.
    """
    values = np.asarray(p_values, dtype=float)
    adjusted = np.full(values.shape, np.nan, dtype=float)
    finite_mask = np.isfinite(values)
    finite = values[finite_mask]
    if finite.size == 0:
        return adjusted

    order = np.argsort(finite)
    ranked = finite[order]
    multipliers = finite.size - np.arange(finite.size)
    step_down = np.maximum.accumulate(ranked * multipliers)
    step_down = np.clip(step_down, 0.0, 1.0)

    restored = np.empty_like(step_down)
    restored[order] = step_down
    adjusted[finite_mask] = restored
    return adjusted


def benjamini_hochberg_pvalues(p_values: Sequence[float]) -> np.ndarray:
    """Returns Benjamini-Hochberg false-discovery-rate adjusted p-values."""
    values = np.asarray(p_values, dtype=float)
    adjusted = np.full(values.shape, np.nan, dtype=float)
    finite_mask = np.isfinite(values)
    finite = values[finite_mask]
    if finite.size == 0:
        return adjusted

    order = np.argsort(finite)
    ranked = finite[order]
    raw = ranked * finite.size / np.arange(1, finite.size + 1)
    monotone = np.minimum.accumulate(raw[::-1])[::-1]
    monotone = np.clip(monotone, 0.0, 1.0)

    restored = np.empty_like(monotone)
    restored[order] = monotone
    adjusted[finite_mask] = restored
    return adjusted


def km_survival_at_horizon(
    durations: Sequence[float],
    events: Sequence[int],
    horizon: float,
) -> float:
    """Returns the observed Kaplan-Meier survival estimate at one horizon."""
    kmf = KaplanMeierFitter()
    kmf.fit(np.asarray(durations, dtype=float), event_observed=np.asarray(events, dtype=int))
    predicted = kmf.predict(float(horizon))
    return float(np.asarray(predicted).reshape(-1)[0])


def ipcw_brier_score(
    durations: Sequence[float],
    events: Sequence[int],
    predicted_survival: Sequence[float],
    horizon: float,
) -> float:
    """Calculates an inverse-probability-of-censoring weighted Brier score.

    Why IPCW is needed:
    A subject censored before the evaluation horizon has an unknown event status at that horizon.
    The censoring distribution is therefore estimated with a Kaplan-Meier model and used to weight
    observable cases. This is more defensible than treating censored-before-horizon records as known
    non-events or silently discarding them.
    """
    duration = np.asarray(durations, dtype=float)
    event = np.asarray(events, dtype=int)
    prediction = np.asarray(predicted_survival, dtype=float)

    censor_km = KaplanMeierFitter()
    censor_km.fit(duration, event_observed=1 - event)

    epsilon = 1e-8
    g_horizon = max(float(np.asarray(censor_km.predict(horizon)).reshape(-1)[0]), epsilon)
    contributions = np.zeros(len(duration), dtype=float)

    event_before = (duration <= horizon) & (event == 1)
    still_observed = duration > horizon

    for index in np.where(event_before)[0]:
        time_before_event = max(float(np.nextafter(duration[index], -np.inf)), 0.0)
        g_event = max(
            float(np.asarray(censor_km.predict(time_before_event)).reshape(-1)[0]),
            epsilon,
        )
        contributions[index] = (prediction[index] ** 2) / g_event

    contributions[still_observed] = ((1.0 - prediction[still_observed]) ** 2) / g_horizon
    # Subjects censored before the horizon receive zero contribution because their horizon status is
    # not observable under the IPCW formulation.
    return float(np.mean(contributions))


def modeled_direction(
    hazard_reduction: float,
    survival_uplift: float,
    tolerance: float,
) -> str:
    """Classifies scenario movement using both hazard and survival evidence."""
    if abs(hazard_reduction) <= tolerance and abs(survival_uplift) <= tolerance:
        return "NEUTRAL"
    if hazard_reduction > tolerance and survival_uplift > tolerance:
        return "IMPROVED"
    if hazard_reduction < -tolerance and survival_uplift < -tolerance:
        return "ADVERSE"
    return "MIXED"


def derive_acceptance_status(checks: pd.DataFrame) -> str:
    """Maps detailed check results into a three-level run conclusion."""
    statuses = set(checks.get("Status", pd.Series(dtype=str)).astype(str))
    if "FAIL" in statuses:
        return "FAIL"
    if "REVIEW" in statuses:
        return "PASS_WITH_REVIEW"
    return "PASS"


# ==================================================================================================
# DETERMINISTIC SYNTHETIC RETENTION DEMONSTRATION
# ==================================================================================================

def generate_synthetic_retention_data(
    n_records: int = 7_500,
    seed: int = 42,
    max_followup_months: float = 24.0,
) -> pd.DataFrame:
    """Generates deterministic synthetic customer-lifecycle data.

    Design principles:
    - No real customer data or PII.
    - Independent censoring is generated separately from event timing.
    - Event timing follows a proportional-hazards-style exponential process whose relative hazard
      depends on engagement, adoption, service burden, tenure, and renewal characteristics.
    - Relationships are directional, not deterministic. Mixed-signal records remain present.
    - Engineered interaction and squared terms are included so dependency-safe simulation can be
      demonstrated and independently audited.

    The resulting dataset is a modeling demonstration. Its distributions and coefficients are not
    institution-calibrated retention estimates.
    """
    if n_records < 500:
        raise ValueError("n_records must be at least 500 for stable segmentation and survival modeling.")
    if max_followup_months <= 1:
        raise ValueError("max_followup_months must be greater than 1.")

    rng = np.random.default_rng(seed)

    customer_id = np.array([f"CUST-{i:07d}" for i in range(1, n_records + 1)])
    enterprise_flag = rng.binomial(1, 0.28, n_records)

    # Latent relationship health introduces realistic correlation without becoming a model feature.
    latent_relationship_health = rng.normal(0.0, 1.0, n_records)

    onboarding_completion = np.clip(
        rng.beta(4.5, 2.2, n_records)
        + 0.05 * enterprise_flag
        + 0.04 * latent_relationship_health,
        0.05,
        1.00,
    )

    monthly_active_days = np.clip(
        rng.normal(
            8.0
            + 13.0 * onboarding_completion
            + 1.5 * enterprise_flag
            + 1.2 * latent_relationship_health,
            3.5,
            n_records,
        ),
        0.0,
        30.0,
    )

    product_adoption_count = np.clip(
        np.rint(
            rng.poisson(
                0.7
                + 2.3 * onboarding_completion
                + 0.08 * monthly_active_days
                + 0.4 * enterprise_flag
            )
        ),
        1,
        8,
    ).astype(int)

    support_tickets_90d = np.clip(
        rng.poisson(
            0.8
            + 1.8 * (1.0 - onboarding_completion)
            + 0.35 * np.maximum(0.0, -latent_relationship_health)
        ),
        0,
        12,
    ).astype(int)

    support_resolution_days = np.clip(
        rng.gamma(shape=2.0, scale=1.2, size=n_records)
        + 0.45 * support_tickets_90d
        - 0.8 * enterprise_flag,
        0.25,
        20.0,
    )

    service_incidents_90d = np.clip(
        rng.poisson(
            0.25
            + 0.18 * support_tickets_90d
            + 0.04 * support_resolution_days
        ),
        0,
        8,
    ).astype(int)

    tenure_at_start = np.clip(
        rng.gamma(shape=2.2, scale=5.5, size=n_records) + 3.0 * enterprise_flag,
        0.5,
        60.0,
    )

    renewal_logit = (
        -0.4
        + 1.6 * onboarding_completion
        + 0.08 * monthly_active_days
        + 0.25 * product_adoption_count
        - 0.08 * support_resolution_days
        + 0.4 * enterprise_flag
    )
    auto_renew_probability = 1.0 / (1.0 + np.exp(-renewal_logit))
    auto_renew_flag = rng.binomial(1, np.clip(auto_renew_probability, 0.05, 0.95))

    # Engineered features follow the framework's strict naming rules.
    monthly_active_days_x_product_adoption_count = (
        monthly_active_days * product_adoption_count
    )
    support_tickets_90d_x_support_resolution_days = (
        support_tickets_90d * support_resolution_days
    )
    onboarding_completion_sq = onboarding_completion ** 2

    # The latent log-hazard is constructed from centered / standardized components. This makes the
    # event process directionally coherent while preventing one raw-scale feature from dominating.
    log_relative_hazard = (
        -0.65 * zscore(onboarding_completion)
        -0.35 * zscore(monthly_active_days)
        -0.25 * zscore(product_adoption_count.astype(float))
        +0.20 * zscore(support_tickets_90d.astype(float))
        +0.30 * zscore(support_resolution_days)
        +0.25 * zscore(service_incidents_90d.astype(float))
        -0.15 * zscore(tenure_at_start)
        -0.25 * zscore(auto_renew_flag.astype(float))
        -0.12 * zscore(enterprise_flag.astype(float))
        -0.15 * zscore(monthly_active_days_x_product_adoption_count)
        +0.18 * zscore(support_tickets_90d_x_support_resolution_days)
        -0.10 * zscore(onboarding_completion_sq)
        + rng.normal(0.0, 0.12, n_records)
    )

    # A constant baseline hazard combined with exp(log_relative_hazard) produces data consistent
    # with the proportional-hazards structure used by CoxPH. Censoring is generated independently.
    baseline_monthly_hazard = 0.024
    individual_hazard = baseline_monthly_hazard * np.exp(log_relative_hazard)
    event_time = rng.exponential(scale=1.0 / individual_hazard)
    censor_time = rng.uniform(8.0, max_followup_months, n_records)

    attrition_event = (event_time <= censor_time).astype(int)
    months_observed = np.minimum(event_time, censor_time)
    months_observed = np.clip(months_observed, 0.10, max_followup_months)

    data = pd.DataFrame(
        {
            "customer_id": customer_id,
            "months_observed": months_observed.round(4),
            "attrition_event": attrition_event,
            "onboarding_completion": onboarding_completion.round(6),
            "monthly_active_days": monthly_active_days.round(4),
            "product_adoption_count": product_adoption_count,
            "support_tickets_90d": support_tickets_90d,
            "support_resolution_days": support_resolution_days.round(4),
            "service_incidents_90d": service_incidents_90d,
            "tenure_at_start": tenure_at_start.round(4),
            "auto_renew_flag": auto_renew_flag,
            "enterprise_flag": enterprise_flag,
            # Engineered terms are recalculated below from the stored base-feature values.
            # This avoids false validation differences caused solely by independent rounding.
            "monthly_active_days_x_product_adoption_count": 0.0,
            "support_tickets_90d_x_support_resolution_days": 0.0,
            "onboarding_completion_sq": 0.0,
        }
    )

    # Rebuild engineered terms from the exact persisted base columns. This is the same dependency
    # rule the scenario engine applies after interventions and guarantees a clean baseline audit.
    data["monthly_active_days_x_product_adoption_count"] = (
        data["monthly_active_days"] * data["product_adoption_count"]
    )
    data["support_tickets_90d_x_support_resolution_days"] = (
        data["support_tickets_90d"] * data["support_resolution_days"]
    )
    data["onboarding_completion_sq"] = data["onboarding_completion"] ** 2

    return data


# ==================================================================================================
# CONFIGURATION BUILDERS
# ==================================================================================================

def build_default_retention_config(output_root: Path = Path("outputs")) -> FrameworkConfig:
    """Returns the public synthetic-retention demonstration contract."""
    return FrameworkConfig(
        project_name="Survival Strategy Framework - Synthetic Retention Demo",
        duration_col="months_observed",
        event_col="attrition_event",
        id_col="customer_id",
        time_unit_label="months",
        model_features=(
            "onboarding_completion",
            "monthly_active_days",
            "product_adoption_count",
            "support_tickets_90d",
            "support_resolution_days",
            "service_incidents_90d",
            "tenure_at_start",
            "auto_renew_flag",
            "enterprise_flag",
            "monthly_active_days_x_product_adoption_count",
            "support_tickets_90d_x_support_resolution_days",
            "onboarding_completion_sq",
        ),
        segmentation_features=(
            "onboarding_completion",
            "monthly_active_days",
            "product_adoption_count",
            "support_tickets_90d",
            "support_resolution_days",
            "service_incidents_90d",
            "tenure_at_start",
        ),
        profile_features=(
            "onboarding_completion",
            "monthly_active_days",
            "product_adoption_count",
            "support_tickets_90d",
            "support_resolution_days",
            "service_incidents_90d",
            "tenure_at_start",
            "auto_renew_flag",
        ),
        n_clusters=3,
        risk_quantile=0.75,
        penalizer=0.10,
        l1_ratio=0.0,
        random_state=42,
        standardize_model_features=True,
        cross_validation_folds=5,
        cross_validation_shuffle=True,
        calibration_horizons=(6.0, 12.0, 18.0, 24.0),
        calibration_groups=10,
        calibration_review_threshold=0.10,
        horizon_extrapolation_tolerance_ratio=0.01,
        show_km_confidence_intervals=True,
        run_ph_test=True,
        ph_test_alpha=0.05,
        ph_strata_max_unique=10,
        persona_silhouette_sample_size=2_000,
        persona_silhouette_review_threshold=0.10,
        persona_stability_repeats=5,
        verify_reproducibility=True,
        evaluation_horizon=12.0,
        timeline_points=121,
        neutral_scenario_tolerance=1e-10,
        output_root=output_root,
        save_input_snapshot=True,
    )


def build_default_retention_scenarios() -> tuple[ScenarioDefinition, ...]:
    """Defines control, favorable, combined, and adverse same-cohort sensitivities.

    The no-change control verifies that the scenario plumbing itself does not create movement. The
    adverse stress proves the engine responds in both directions. The favorable scenarios remain
    modeled sensitivities, not causal intervention-effect estimates.
    """
    return (
        ScenarioDefinition(
            name="No-Change Control",
            description=(
                "Apply no feature changes to the baseline target cohort. Expected hazard and "
                "survival movement should be zero within numerical tolerance."
            ),
            changes=(),
            scenario_type="CONTROL",
            expected_direction="NEUTRAL",
        ),
        ScenarioDefinition(
            name="Onboarding Completion Improvement",
            description=(
                "Increase onboarding completion by 10 percentage points, bounded at 100%, for the "
                "baseline top-quartile risk cohort."
            ),
            changes=(
                FeatureChange(
                    feature="onboarding_completion",
                    operation="add",
                    value=0.10,
                    lower_bound=0.0,
                    upper_bound=1.0,
                    round_digits=6,
                ),
            ),
            scenario_type="IMPROVEMENT",
            expected_direction="IMPROVED",
        ),
        ScenarioDefinition(
            name="Support Resolution Improvement",
            description=(
                "Reduce support resolution time by 25%, with a 0.25-day floor, for the baseline "
                "top-quartile risk cohort."
            ),
            changes=(
                FeatureChange(
                    feature="support_resolution_days",
                    operation="multiply",
                    value=0.75,
                    lower_bound=0.25,
                    upper_bound=20.0,
                    round_digits=4,
                ),
            ),
            scenario_type="IMPROVEMENT",
            expected_direction="IMPROVED",
        ),
        ScenarioDefinition(
            name="Product Adoption Expansion",
            description=(
                "Increase adopted products by one, bounded at eight, for the baseline top-quartile "
                "risk cohort."
            ),
            changes=(
                FeatureChange(
                    feature="product_adoption_count",
                    operation="add",
                    value=1.0,
                    lower_bound=1.0,
                    upper_bound=8.0,
                    round_digits=0,
                ),
            ),
            scenario_type="IMPROVEMENT",
            expected_direction="IMPROVED",
        ),
        ScenarioDefinition(
            name="Combined Retention Strategy",
            description=(
                "Apply the onboarding, support-resolution, and product-adoption assumptions "
                "together to the same baseline top-quartile risk cohort."
            ),
            changes=(
                FeatureChange(
                    feature="onboarding_completion",
                    operation="add",
                    value=0.10,
                    lower_bound=0.0,
                    upper_bound=1.0,
                    round_digits=6,
                ),
                FeatureChange(
                    feature="support_resolution_days",
                    operation="multiply",
                    value=0.75,
                    lower_bound=0.25,
                    upper_bound=20.0,
                    round_digits=4,
                ),
                FeatureChange(
                    feature="product_adoption_count",
                    operation="add",
                    value=1.0,
                    lower_bound=1.0,
                    upper_bound=8.0,
                    round_digits=0,
                ),
            ),
            scenario_type="IMPROVEMENT",
            expected_direction="IMPROVED",
        ),
        ScenarioDefinition(
            name="Service Friction Stress",
            description=(
                "Increase support resolution time by 25% and add one service incident, within "
                "governed bounds, to verify adverse bidirectional sensitivity for the same cohort."
            ),
            changes=(
                FeatureChange(
                    feature="support_resolution_days",
                    operation="multiply",
                    value=1.25,
                    lower_bound=0.25,
                    upper_bound=20.0,
                    round_digits=4,
                ),
                FeatureChange(
                    feature="service_incidents_90d",
                    operation="add",
                    value=1.0,
                    lower_bound=0.0,
                    upper_bound=8.0,
                    round_digits=0,
                ),
            ),
            scenario_type="STRESS",
            expected_direction="ADVERSE",
        ),
    )

# ==================================================================================================
# INPUT CONTRACT AND VALIDATION GATEKEEPER
# ==================================================================================================

def engineered_feature_sources(feature: str) -> tuple[str, ...] | None:
    """Returns source columns implied by the strict engineered-feature naming rules."""
    if "_x_" in feature:
        return tuple(part for part in feature.split("_x_") if part)
    if feature.endswith("_sq"):
        return (feature[:-3],)
    return None


def validate_survival_input(
    df: pd.DataFrame,
    config: FrameworkConfig,
    scenarios: Sequence[ScenarioDefinition],
) -> ValidationResult:
    """Validates the complete analytical contract before segmentation or model fitting.

    The gatekeeper deliberately separates hard failures from review-level warnings. A reusable
    framework should fail fast on invalid IDs, durations, event coding, feature types, stale
    engineered terms, or unsafe scenarios. It should record—but not hide—issues such as low event
    density or extrapolative evaluation horizons.
    """
    result = ValidationResult()

    if not isinstance(df, pd.DataFrame):
        result.add("ERROR", "Dataset", "DATAFRAME", "Input must be a pandas DataFrame.")
        return result
    if df.empty:
        result.add("ERROR", "Dataset", "ROW_COUNT", "Input dataset contains no rows.")
        return result

    required_columns = {
        config.id_col,
        config.duration_col,
        config.event_col,
        *config.model_features,
        *config.segmentation_features,
    }
    missing_columns = sorted(required_columns.difference(df.columns))
    for column in missing_columns:
        result.add("ERROR", "Required Column", column, "Required field is missing from the input.")

    # Configuration checks can run even when input fields are missing.
    if not config.project_name.strip():
        result.add("ERROR", "Configuration", "project_name", "Project name may not be blank.")
    if not config.time_unit_label.strip():
        result.add("ERROR", "Configuration", "time_unit_label", "Time-unit label may not be blank.")
    if not 0.0 < config.risk_quantile < 1.0:
        result.add("ERROR", "Configuration", "risk_quantile", "Must be between 0 and 1.")
    if config.n_clusters < 2:
        result.add("ERROR", "Configuration", "n_clusters", "Must be at least 2.")
    if config.n_clusters >= len(df):
        result.add("ERROR", "Configuration", "n_clusters", "Must be smaller than row count.")
    if config.penalizer < 0:
        result.add("ERROR", "Configuration", "penalizer", "Must be nonnegative.")
    if not 0.0 <= config.l1_ratio <= 1.0:
        result.add("ERROR", "Configuration", "l1_ratio", "Must be between 0 and 1.")
    if config.cross_validation_folds < 2:
        result.add("ERROR", "Configuration", "cross_validation_folds", "Must be at least 2.")
    if config.calibration_groups < 2:
        result.add("ERROR", "Configuration", "calibration_groups", "Must be at least 2.")
    if config.horizon_extrapolation_tolerance_ratio < 0:
        result.add(
            "ERROR",
            "Configuration",
            "horizon_extrapolation_tolerance_ratio",
            "Must be nonnegative.",
        )
    if not config.calibration_horizons:
        result.add("ERROR", "Configuration", "calibration_horizons", "At least one horizon is required.")
    if any(horizon <= 0 for horizon in config.calibration_horizons):
        result.add("ERROR", "Configuration", "calibration_horizons", "All horizons must be positive.")
    if config.evaluation_horizon <= 0:
        result.add("ERROR", "Configuration", "evaluation_horizon", "Must be positive.")
    if config.timeline_points < 20:
        result.add("ERROR", "Configuration", "timeline_points", "Must be at least 20.")
    if config.persona_stability_repeats < 1:
        result.add("ERROR", "Configuration", "persona_stability_repeats", "Must be at least 1.")
    if config.neutral_scenario_tolerance <= 0:
        result.add("ERROR", "Configuration", "neutral_scenario_tolerance", "Must be positive.")

    scenario_names = [scenario.name for scenario in scenarios]
    if len(set(scenario_names)) != len(scenario_names):
        result.add("ERROR", "Scenario", "Scenario_Name", "Scenario names must be unique.")

    valid_scenario_types = {"CONTROL", "IMPROVEMENT", "STRESS", "TECHNICAL"}
    valid_directions = {"NEUTRAL", "IMPROVED", "ADVERSE", "UNSPECIFIED"}
    for scenario in scenarios:
        if scenario.scenario_type not in valid_scenario_types:
            result.add("ERROR", "Scenario", scenario.name, "Unsupported scenario_type.")
        if scenario.expected_direction not in valid_directions:
            result.add("ERROR", "Scenario", scenario.name, "Unsupported expected_direction.")
        if scenario.scenario_type == "CONTROL" and scenario.changes:
            result.add(
                "WARNING",
                "Scenario",
                scenario.name,
                "CONTROL scenario contains feature changes; confirm that this is intentional.",
            )
        if not scenario.changes and scenario.scenario_type != "CONTROL":
            result.add(
                "WARNING",
                "Scenario",
                scenario.name,
                "Scenario contains no feature changes but is not labeled CONTROL.",
            )

    if missing_columns:
        return result

    # Governed ID validation.
    if df[config.id_col].isna().any():
        result.add("ERROR", "Identifier", config.id_col, "ID contains missing values.")
    if df[config.id_col].duplicated().any():
        duplicate_count = int(df[config.id_col].duplicated(keep=False).sum())
        result.add(
            "ERROR",
            "Identifier",
            config.id_col,
            f"ID is not unique; {duplicate_count:,} rows participate in duplicates.",
        )

    # Duration validation.
    duration = pd.to_numeric(df[config.duration_col], errors="coerce")
    if duration.isna().any():
        result.add("ERROR", "Time-to-Event", config.duration_col, "Duration must be numeric and non-missing.")
    else:
        if not np.isfinite(duration.to_numpy()).all():
            result.add("ERROR", "Time-to-Event", config.duration_col, "Duration contains infinite values.")
        if (duration <= 0).any():
            result.add("ERROR", "Time-to-Event", config.duration_col, "Duration must be strictly positive.")
        max_duration = float(duration.max())
        horizon_tolerance = max(
            1e-8,
            abs(max_duration) * config.horizon_extrapolation_tolerance_ratio,
        )
        for horizon in sorted(set((*config.calibration_horizons, config.evaluation_horizon))):
            if horizon > max_duration + horizon_tolerance:
                result.add(
                    "WARNING",
                    "Time-to-Event",
                    f"horizon_{horizon:g}",
                    (
                        "Evaluation horizon materially exceeds the maximum observed duration; "
                        "predicted survival is extrapolative."
                    ),
                )

    # Event validation.
    event = pd.to_numeric(df[config.event_col], errors="coerce")
    if event.isna().any():
        result.add("ERROR", "Event", config.event_col, "Event field must be binary and non-missing.")
    else:
        invalid_events = sorted(set(event.unique()).difference({0, 1}))
        if invalid_events:
            result.add(
                "ERROR",
                "Event",
                config.event_col,
                f"Event field contains values outside {{0, 1}}: {invalid_events[:10]}",
            )
        else:
            event_rate = float(event.mean())
            result.add(
                "INFO",
                "Event",
                config.event_col,
                f"Observed event rate is {event_rate:.2%}; right-censoring rate is {1.0 - event_rate:.2%}.",
            )
            if event_rate < config.low_event_rate_warning:
                result.add("WARNING", "Event", config.event_col, "Event rate is very low; model stability requires review.")
            if event_rate > config.high_event_rate_warning:
                result.add("WARNING", "Event", config.event_col, "Event rate is very high; censoring design requires review.")

            class_counts = event.value_counts()
            if len(class_counts) < 2:
                result.add("ERROR", "Cross-Validation", config.event_col, "Both events and censored records are required.")
            elif config.cross_validation_folds > int(class_counts.min()):
                result.add(
                    "ERROR",
                    "Cross-Validation",
                    "cross_validation_folds",
                    "Fold count exceeds the smaller event/censoring class and cannot be stratified safely.",
                )

    # Feature validation. This release requires numeric / Boolean inputs.
    all_features = list(dict.fromkeys((*config.model_features, *config.segmentation_features)))
    for feature in all_features:
        series = df[feature]
        if not (pd.api.types.is_numeric_dtype(series) or pd.api.types.is_bool_dtype(series)):
            result.add(
                "ERROR",
                "Feature Contract",
                feature,
                "Feature must be numeric or Boolean; automatic categorical encoding is not claimed in this release.",
            )
            continue

        numeric = pd.to_numeric(series, errors="coerce")
        missing_ratio = float(numeric.isna().mean())
        if missing_ratio > 0:
            result.add("ERROR", "Feature Contract", feature, f"Feature contains {missing_ratio:.2%} missing values.")
        if not np.isfinite(numeric.fillna(0).to_numpy()).all():
            result.add("ERROR", "Feature Contract", feature, "Feature contains infinite values.")
        if numeric.nunique(dropna=True) <= 1:
            result.add("ERROR", "Feature Contract", feature, "Feature has zero variance.")
        elif missing_ratio > config.high_missingness_warning:
            result.add("WARNING", "Feature Contract", feature, "Feature missingness exceeds review threshold.")

    # Engineered-feature consistency is checked before any scenario is applied.
    for feature in config.model_features:
        sources = engineered_feature_sources(feature)
        if not sources:
            continue
        missing_sources = [source for source in sources if source not in df.columns]
        if missing_sources:
            result.add(
                "ERROR",
                "Engineered Feature",
                feature,
                f"Source field(s) missing: {missing_sources}",
            )
            continue
        if "_x_" in feature:
            expected = df.loc[:, list(sources)].prod(axis=1)
        else:
            expected = pd.to_numeric(df[sources[0]], errors="coerce") ** 2
        actual = pd.to_numeric(df[feature], errors="coerce")
        difference = np.nanmax(np.abs(actual.to_numpy() - expected.to_numpy()))
        if not np.isfinite(difference) or difference > 1e-8:
            result.add(
                "ERROR",
                "Engineered Feature",
                feature,
                f"Stored values are stale or inconsistent; max absolute difference is {difference}.",
            )

    protected_fields = {config.id_col, config.duration_col, config.event_col}
    model_feature_set = set(config.model_features)
    for scenario in scenarios:
        for change in scenario.changes:
            if change.feature not in df.columns:
                result.add("ERROR", "Scenario", scenario.name, f"Feature '{change.feature}' does not exist.")
            if change.feature in protected_fields:
                result.add("ERROR", "Scenario", scenario.name, f"Protected field '{change.feature}' may not be changed.")
            if change.feature in model_feature_set and engineered_feature_sources(change.feature):
                result.add(
                    "ERROR",
                    "Scenario",
                    scenario.name,
                    f"Engineered field '{change.feature}' may not be changed directly; change its base feature(s).",
                )
            if change.operation not in {"multiply", "add", "replace"}:
                result.add("ERROR", "Scenario", scenario.name, f"Unsupported operation '{change.operation}'.")
            if (
                change.lower_bound is not None
                and change.upper_bound is not None
                and change.lower_bound > change.upper_bound
            ):
                result.add("ERROR", "Scenario", scenario.name, "Lower bound exceeds upper bound.")

    return result

# ==================================================================================================
# PERSONA BRANCH — UNSUPERVISED SEGMENTATION
# ==================================================================================================

def derive_persona_names(
    standardized_centroids: np.ndarray,
    features: Sequence[str],
) -> dict[int, str]:
    """Assigns descriptive lifecycle labels from standardized persona centroids.

    The rules are deliberately transparent. They do not claim that K-Means discovered causal
    customer types. They translate the strongest centroid pattern into a readable label while
    preserving Segment_ID for exact technical traceability.
    """
    feature_index = {feature: index for index, feature in enumerate(features)}

    def average(row: np.ndarray, names: Sequence[str]) -> float:
        values = [row[feature_index[name]] for name in names if name in feature_index]
        return float(np.mean(values)) if values else 0.0

    labels: dict[int, str] = {}
    used: dict[str, int] = {}
    for segment_id, centroid in enumerate(standardized_centroids):
        engagement = average(
            centroid,
            ("onboarding_completion", "monthly_active_days", "product_adoption_count"),
        )
        friction = average(
            centroid,
            ("support_tickets_90d", "support_resolution_days", "service_incidents_90d"),
        )
        tenure = average(centroid, ("tenure_at_start",))

        if engagement >= 0.35 and friction <= 0.05:
            base = "Engaged & Embedded"
        elif friction >= 0.40:
            base = "Service-Friction Risk"
        elif engagement <= -0.30:
            base = "Low-Engagement Risk"
        elif tenure >= 0.35 and friction <= 0.15:
            base = "Established & Stable"
        else:
            base = "Mixed Lifecycle Profile"

        used[base] = used.get(base, 0) + 1
        labels[segment_id] = base if used[base] == 1 else f"{base} ({used[base]})"
    return labels


def fit_personas(df: pd.DataFrame, config: FrameworkConfig) -> PersonaArtifacts:
    """Fits reproducible K-Means personas independently from the CoxPH model.

    Quality evidence includes sampled silhouette, Davies-Bouldin, Calinski-Harabasz, cluster size
    balance, and multi-seed adjusted Rand index. These measures help future users assess whether the
    persona solution is stable and interpretable rather than accepting cluster labels uncritically.
    """
    features = numeric_frame(df, config.segmentation_features)
    scaler = StandardScaler()
    scaled = scaler.fit_transform(features)

    kmeans = KMeans(
        n_clusters=config.n_clusters,
        random_state=config.random_state,
        n_init=20,
    )
    raw_labels = kmeans.fit_predict(scaled)
    persona_name_map = derive_persona_names(kmeans.cluster_centers_, config.segmentation_features)

    scored = df.copy()
    scored["Segment_ID"] = raw_labels.astype(int)
    scored["Persona"] = scored["Segment_ID"].map(persona_name_map)

    grouped = scored.groupby(["Segment_ID", "Persona"], sort=True)
    profiles = grouped.agg(
        Record_Count=(config.id_col, "count"),
        Event_Rate=(config.event_col, "mean"),
        Median_Duration=(config.duration_col, "median"),
    ).reset_index()

    profile_features = config.profile_features or config.segmentation_features
    feature_means = grouped[list(profile_features)].mean().reset_index()
    profiles = profiles.merge(feature_means, on=["Segment_ID", "Persona"], how="left")

    # Silhouette can become memory-intensive on large datasets, so the framework uses a deterministic
    # sample while calculating the other metrics on the full scaled matrix.
    silhouette_sample = min(config.persona_silhouette_sample_size, len(df))
    silhouette = float(
        silhouette_score(
            scaled,
            raw_labels,
            sample_size=silhouette_sample if silhouette_sample < len(df) else None,
            random_state=config.random_state,
        )
    )
    davies_bouldin = float(davies_bouldin_score(scaled, raw_labels))
    calinski_harabasz = float(calinski_harabasz_score(scaled, raw_labels))

    cluster_counts = pd.Series(raw_labels).value_counts().sort_index()
    smallest_cluster_rate = float(cluster_counts.min() / len(df))
    largest_cluster_rate = float(cluster_counts.max() / len(df))

    stability_rows: list[dict[str, Any]] = []
    for repeat in range(1, config.persona_stability_repeats + 1):
        alternate_seed = config.random_state + repeat
        alternate = KMeans(
            n_clusters=config.n_clusters,
            random_state=alternate_seed,
            n_init=20,
        ).fit_predict(scaled)
        stability_rows.append(
            {
                "Repeat": repeat,
                "Random_State": alternate_seed,
                "Adjusted_Rand_Index_vs_Baseline": float(
                    adjusted_rand_score(raw_labels, alternate)
                ),
            }
        )
    stability_results = pd.DataFrame(stability_rows)

    quality_metrics = pd.DataFrame(
        [
            {
                "Rows": int(len(df)),
                "Clusters": int(config.n_clusters),
                "Inertia": float(kmeans.inertia_),
                "Silhouette_Score": silhouette,
                "Silhouette_Sample_Size": int(silhouette_sample),
                "Davies_Bouldin_Index": davies_bouldin,
                "Calinski_Harabasz_Index": calinski_harabasz,
                "Smallest_Cluster_Rate": smallest_cluster_rate,
                "Largest_Cluster_Rate": largest_cluster_rate,
                "Stability_Mean_ARI": float(
                    stability_results["Adjusted_Rand_Index_vs_Baseline"].mean()
                ),
                "Stability_Min_ARI": float(
                    stability_results["Adjusted_Rand_Index_vs_Baseline"].min()
                ),
                "Silhouette_Review_Threshold": config.persona_silhouette_review_threshold,
                "Persona_Quality_Status": (
                    "PASS"
                    if silhouette >= config.persona_silhouette_review_threshold
                    else "REVIEW"
                ),
            }
        ]
    )

    return PersonaArtifacts(
        scored_data=scored,
        scaler=scaler,
        kmeans=kmeans,
        profiles=profiles,
        quality_metrics=quality_metrics,
        stability_results=stability_results,
        persona_name_map=persona_name_map,
    )


# ==================================================================================================
# COXPH BRANCH — MODEL FITTING, CROSS-VALIDATION, CALIBRATION, AND PH REVIEW
# ==================================================================================================

def transform_model_features(
    df: pd.DataFrame,
    config: FrameworkConfig,
    scaler: StandardScaler | None,
) -> pd.DataFrame:
    raw = numeric_frame(df, config.model_features)
    if scaler is None:
        return raw
    transformed = scaler.transform(raw)
    return pd.DataFrame(transformed, columns=list(config.model_features), index=df.index)


def _model_summary_frame(cph: CoxPHFitter, standardized: bool) -> pd.DataFrame:
    summary = cph.summary.reset_index()
    first_column = summary.columns[0]
    summary = summary.rename(columns={first_column: "Feature"})

    rename_map = {
        "coef": "Coefficient",
        "exp(coef)": "Hazard_Ratio",
        "se(coef)": "Standard_Error",
        "coef lower 95%": "Coefficient_CI_Lower",
        "coef upper 95%": "Coefficient_CI_Upper",
        "exp(coef) lower 95%": "Hazard_Ratio_CI_Lower",
        "exp(coef) upper 95%": "Hazard_Ratio_CI_Upper",
        "z": "Z_Statistic",
        "p": "P_Value",
        "-log2(p)": "Negative_Log2_P",
    }
    summary = summary.rename(columns=rename_map)

    desired = [
        "Feature",
        "Coefficient",
        "Hazard_Ratio",
        "Standard_Error",
        "Coefficient_CI_Lower",
        "Coefficient_CI_Upper",
        "Hazard_Ratio_CI_Lower",
        "Hazard_Ratio_CI_Upper",
        "Z_Statistic",
        "P_Value",
        "Negative_Log2_P",
    ]
    available = [column for column in desired if column in summary.columns]
    summary = summary.loc[:, available].copy()
    summary["Interpretation_Unit"] = "Per 1 SD" if standardized else "Per raw feature unit"
    return summary


def cross_validate_cox_model(
    df: pd.DataFrame,
    config: FrameworkConfig,
) -> tuple[pd.DataFrame, pd.DataFrame]:
    """Runs deterministic stratified K-fold validation and returns out-of-fold predictions.

    Every fold fits its own scaler on training data only. This avoids leaking holdout-distribution
    information into the validation fold. The out-of-fold survival predictions are subsequently used
    for calibration evidence and IPCW Brier scores.
    """
    splitter = StratifiedKFold(
        n_splits=config.cross_validation_folds,
        shuffle=config.cross_validation_shuffle,
        random_state=config.random_state if config.cross_validation_shuffle else None,
    )

    event_vector = pd.to_numeric(df[config.event_col], errors="raise").astype(int).to_numpy()
    fold_rows: list[dict[str, Any]] = []
    prediction_frames: list[pd.DataFrame] = []
    horizons = np.asarray(sorted(set(config.calibration_horizons)), dtype=float)

    for fold_number, (train_positions, validation_positions) in enumerate(
        splitter.split(np.zeros(len(df)), event_vector),
        start=1,
    ):
        train_df = df.iloc[train_positions].copy()
        validation_df = df.iloc[validation_positions].copy()
        try:
            train_raw = numeric_frame(train_df, config.model_features)
            validation_raw = numeric_frame(validation_df, config.model_features)
            fold_scaler: StandardScaler | None = None

            if config.standardize_model_features:
                fold_scaler = StandardScaler()
                train_x = pd.DataFrame(
                    fold_scaler.fit_transform(train_raw),
                    columns=list(config.model_features),
                    index=train_df.index,
                )
                validation_x = pd.DataFrame(
                    fold_scaler.transform(validation_raw),
                    columns=list(config.model_features),
                    index=validation_df.index,
                )
            else:
                train_x = train_raw
                validation_x = validation_raw

            train_frame = train_x.copy()
            train_frame[config.duration_col] = pd.to_numeric(
                train_df[config.duration_col], errors="raise"
            )
            train_frame[config.event_col] = pd.to_numeric(
                train_df[config.event_col], errors="raise"
            ).astype(int)

            fold_model = CoxPHFitter(
                penalizer=config.penalizer,
                l1_ratio=config.l1_ratio,
            )
            fold_model.fit(
                train_frame,
                duration_col=config.duration_col,
                event_col=config.event_col,
                show_progress=False,
            )

            validation_hazard = np.asarray(
                fold_model.predict_partial_hazard(validation_x)
            ).reshape(-1)
            validation_duration = pd.to_numeric(
                validation_df[config.duration_col], errors="raise"
            ).to_numpy(dtype=float)
            validation_event = pd.to_numeric(
                validation_df[config.event_col], errors="raise"
            ).to_numpy(dtype=int)
            fold_concordance = float(
                concordance_index(
                    validation_duration,
                    -validation_hazard,
                    validation_event,
                )
            )

            survival = fold_model.predict_survival_function(
                validation_x,
                times=horizons,
            )
            survival_values = survival.to_numpy(dtype=float).T

            prediction_frame = pd.DataFrame(
                {
                    "Original_Index": validation_df.index.to_numpy(),
                    config.id_col: validation_df[config.id_col].to_numpy(),
                    config.duration_col: validation_duration,
                    config.event_col: validation_event,
                    "CV_Fold": fold_number,
                    "OOF_Partial_Hazard": validation_hazard,
                }
            )
            for horizon_index, horizon in enumerate(horizons):
                prediction_frame[
                    f"OOF_Survival_At_{horizon_token(horizon)}"
                ] = survival_values[:, horizon_index]
            prediction_frames.append(prediction_frame)

            fold_rows.append(
                {
                    "Fold": fold_number,
                    "Train_Rows": int(len(train_df)),
                    "Validation_Rows": int(len(validation_df)),
                    "Validation_Events": int(validation_event.sum()),
                    "Validation_Event_Rate": float(validation_event.mean()),
                    "Concordance_Index": fold_concordance,
                    "Status": "PASS",
                    "Error": "",
                }
            )
        except Exception as exc:
            fold_rows.append(
                {
                    "Fold": fold_number,
                    "Train_Rows": int(len(train_df)),
                    "Validation_Rows": int(len(validation_df)),
                    "Validation_Events": int(validation_df[config.event_col].sum()),
                    "Validation_Event_Rate": float(validation_df[config.event_col].mean()),
                    "Concordance_Index": np.nan,
                    "Status": "FAIL",
                    "Error": str(exc),
                }
            )

    fold_results = pd.DataFrame(fold_rows)
    predictions = (
        pd.concat(prediction_frames, ignore_index=True)
        if prediction_frames
        else pd.DataFrame()
    )
    if not predictions.empty:
        predictions = predictions.sort_values("Original_Index").reset_index(drop=True)
    return fold_results, predictions


def build_calibration_evidence(
    cross_validated_predictions: pd.DataFrame,
    config: FrameworkConfig,
) -> tuple[pd.DataFrame, pd.DataFrame]:
    """Builds horizon-specific out-of-fold calibration and IPCW Brier evidence."""
    if cross_validated_predictions.empty:
        error = pd.DataFrame(
            [{
                "Horizon": np.nan,
                "Time_Unit": config.time_unit_label,
                "Calibration_Status": "NOT_AVAILABLE",
                "Error": "No out-of-fold predictions were available.",
            }]
        )
        return error, pd.DataFrame()

    predictions = cross_validated_predictions.copy()
    ranked = predictions["OOF_Partial_Hazard"].rank(method="first")
    group_count = min(config.calibration_groups, len(predictions))
    group_codes = pd.qcut(ranked, q=group_count, labels=False, duplicates="drop")
    predictions["Calibration_Group_Order"] = group_codes.astype(int) + 1
    predictions["Calibration_Risk_Group"] = predictions["Calibration_Group_Order"].map(
        lambda value: (
            f"G{int(value):02d} Lowest Modeled Hazard"
            if int(value) == 1
            else (
                f"G{int(value):02d} Highest Modeled Hazard"
                if int(value) == group_count
                else f"G{int(value):02d}"
            )
        )
    )

    metric_rows: list[dict[str, Any]] = []
    group_rows: list[dict[str, Any]] = []
    brier_values: list[tuple[float, float]] = []

    for horizon in sorted(set(config.calibration_horizons)):
        survival_column = f"OOF_Survival_At_{horizon_token(horizon)}"
        predicted = pd.to_numeric(predictions[survival_column], errors="raise")
        observed = km_survival_at_horizon(
            predictions[config.duration_col],
            predictions[config.event_col],
            horizon,
        )
        brier = ipcw_brier_score(
            predictions[config.duration_col],
            predictions[config.event_col],
            predicted,
            horizon,
        )
        brier_values.append((horizon, brier))

        duration_values = pd.to_numeric(
            predictions[config.duration_col], errors="raise"
        )
        event_values = pd.to_numeric(
            predictions[config.event_col], errors="raise"
        ).astype(int)
        metric_rows.append(
            {
                "Horizon": float(horizon),
                "Time_Unit": config.time_unit_label,
                "At_Risk_At_Horizon": int((duration_values >= horizon).sum()),
                "Observed_Events_By_Horizon": int(
                    ((duration_values <= horizon) & event_values.eq(1)).sum()
                ),
                "Censored_By_Horizon": int(
                    ((duration_values <= horizon) & event_values.eq(0)).sum()
                ),
                "Overall_Predicted_Survival": float(predicted.mean()),
                "Overall_Observed_KM_Survival": float(observed),
                "Calibration_Difference": float(predicted.mean() - observed),
                "Absolute_Calibration_Error": float(abs(predicted.mean() - observed)),
                "IPCW_Brier_Score": float(brier),
                "Brier_Score_Method": "IPCW",
            }
        )

        for (group_order, group_name), group in predictions.groupby(
            ["Calibration_Group_Order", "Calibration_Risk_Group"],
            sort=True,
        ):
            group_predicted = pd.to_numeric(group[survival_column], errors="raise")
            group_observed = km_survival_at_horizon(
                group[config.duration_col],
                group[config.event_col],
                horizon,
            )
            group_duration = pd.to_numeric(group[config.duration_col], errors="raise")
            group_event = pd.to_numeric(group[config.event_col], errors="raise").astype(int)
            group_rows.append(
                {
                    "Horizon": float(horizon),
                    "Time_Unit": config.time_unit_label,
                    "Calibration_Group_Order": int(group_order),
                    "Calibration_Risk_Group": str(group_name),
                    "Record_Count": int(len(group)),
                    "Event_Count": int(group_event.sum()),
                    "At_Risk_At_Horizon": int((group_duration >= horizon).sum()),
                    "Observed_Events_By_Horizon": int(
                        ((group_duration <= horizon) & group_event.eq(1)).sum()
                    ),
                    "Censored_By_Horizon": int(
                        ((group_duration <= horizon) & group_event.eq(0)).sum()
                    ),
                    "Predicted_Survival": float(group_predicted.mean()),
                    "Observed_KM_Survival": float(group_observed),
                    "Calibration_Difference": float(group_predicted.mean() - group_observed),
                    "Absolute_Calibration_Error": float(
                        abs(group_predicted.mean() - group_observed)
                    ),
                }
            )

    metrics = pd.DataFrame(metric_rows)
    by_group = pd.DataFrame(group_rows)

    horizon_array = np.asarray([item[0] for item in brier_values], dtype=float)
    brier_array = np.asarray([item[1] for item in brier_values], dtype=float)
    if len(horizon_array) > 1 and float(horizon_array.max() - horizon_array.min()) > 0:
        integration = getattr(np, "trapezoid", None)
        if integration is None:  # Compatibility with older NumPy releases.
            integration = np.trapz
        integrated_brier = float(
            integration(brier_array, horizon_array)
            / (horizon_array.max() - horizon_array.min())
        )
    else:
        integrated_brier = float(brier_array[0])

    max_group_error = (
        float(by_group["Absolute_Calibration_Error"].max())
        if not by_group.empty
        else np.nan
    )
    metrics["Integrated_Brier_Score"] = integrated_brier
    metrics["Integrated_Brier_Score_Scope"] = (
        "Trapezoidal average across configured horizons; not a full continuous-time IBS."
    )
    metrics["Integration_Horizon_Min"] = float(horizon_array.min())
    metrics["Integration_Horizon_Max"] = float(horizon_array.max())
    metrics["Max_Group_Absolute_Calibration_Error"] = max_group_error
    metrics["Calibration_Review_Threshold"] = config.calibration_review_threshold
    metrics["Calibration_Status"] = np.where(
        max_group_error <= config.calibration_review_threshold,
        "PASS",
        "REVIEW",
    )
    return metrics, by_group


def _extract_scaled_schoenfeld_residuals(
    cph: CoxPHFitter,
    training_frame: pd.DataFrame,
    duration_col: str,
    review_features: Sequence[str],
) -> pd.DataFrame:
    if not review_features:
        return pd.DataFrame(
            columns=["Row_Index", "Event_Time", "Feature", "Scaled_Schoenfeld_Residual"]
        )
    try:
        residuals = cph.compute_residuals(training_frame, kind="scaled_schoenfeld")
        rows: list[pd.DataFrame] = []
        for feature in review_features:
            if feature not in residuals.columns:
                continue
            frame = pd.DataFrame(
                {
                    "Row_Index": residuals.index,
                    "Event_Time": training_frame.loc[residuals.index, duration_col].to_numpy(),
                    "Feature": feature,
                    "Scaled_Schoenfeld_Residual": residuals[feature].to_numpy(),
                }
            )
            rows.append(frame)
        return pd.concat(rows, ignore_index=True) if rows else pd.DataFrame()
    except Exception as exc:
        return pd.DataFrame(
            [{
                "Row_Index": np.nan,
                "Event_Time": np.nan,
                "Feature": "ALL",
                "Scaled_Schoenfeld_Residual": np.nan,
                "Error": str(exc),
            }]
        )


def fit_cox_model(df: pd.DataFrame, config: FrameworkConfig) -> CoxModelArtifacts:
    """Fits the full regularized CoxPH model and all model-development evidence layers."""
    raw_features = numeric_frame(df, config.model_features)
    model_scaler: StandardScaler | None = None

    if config.standardize_model_features:
        model_scaler = StandardScaler()
        transformed_features = pd.DataFrame(
            model_scaler.fit_transform(raw_features),
            columns=list(config.model_features),
            index=df.index,
        )
    else:
        transformed_features = raw_features

    training_frame = transformed_features.copy()
    training_frame[config.duration_col] = pd.to_numeric(df[config.duration_col], errors="raise")
    training_frame[config.event_col] = pd.to_numeric(df[config.event_col], errors="raise").astype(int)

    cph = CoxPHFitter(penalizer=config.penalizer, l1_ratio=config.l1_ratio)
    captured_warnings: list[str] = []
    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always")
        cph.fit(
            training_frame,
            duration_col=config.duration_col,
            event_col=config.event_col,
            show_progress=False,
        )
        captured_warnings = [str(item.message) for item in caught]

    model_summary = _model_summary_frame(cph, config.standardize_model_features)
    cv_results, cv_predictions = cross_validate_cox_model(df, config)
    calibration_metrics, calibration_by_group = build_calibration_evidence(
        cv_predictions,
        config,
    )

    if config.run_ph_test:
        try:
            ph_result = proportional_hazard_test(
                cph,
                training_frame,
                time_transform="rank",
            )
            ph_test_results = ph_result.summary.reset_index()
            first_column = ph_test_results.columns[0]
            ph_test_results = ph_test_results.rename(
                columns={
                    first_column: "Feature",
                    "test_statistic": "Test_Statistic",
                    "p": "P_Value",
                    "-log2(p)": "Negative_Log2_P",
                }
            )
            ph_test_results["P_Value_Holm"] = holm_adjusted_pvalues(
                ph_test_results["P_Value"]
            )
            ph_test_results["P_Value_BH"] = benjamini_hochberg_pvalues(
                ph_test_results["P_Value"]
            )
            ph_test_results["Raw_PH_Status"] = np.where(
                ph_test_results["P_Value"] < config.ph_test_alpha,
                "REVIEW",
                "PASS",
            )
            ph_test_results["Holm_PH_Status"] = np.where(
                ph_test_results["P_Value_Holm"] < config.ph_test_alpha,
                "REVIEW",
                "PASS",
            )
            ph_test_results["BH_PH_Status"] = np.where(
                ph_test_results["P_Value_BH"] < config.ph_test_alpha,
                "REVIEW",
                "PASS",
            )
            ph_test_results["PH_Assumption_Status"] = np.select(
                [
                    ph_test_results["Holm_PH_Status"].eq("REVIEW"),
                    ph_test_results["Raw_PH_Status"].eq("REVIEW"),
                ],
                ["REVIEW", "PASS_WITH_RAW_REVIEW"],
                default="PASS",
            )
            ph_test_results["Review_Rationale"] = np.select(
                [
                    ph_test_results["Holm_PH_Status"].eq("REVIEW"),
                    ph_test_results["Raw_PH_Status"].eq("REVIEW"),
                ],
                [
                    "Feature remains below alpha after Holm adjustment; sensitivity review required.",
                    "Raw p-value is below alpha but multiplicity-adjusted tests pass; retain as documented review item.",
                ],
                default="No feature-level PH review item at the configured alpha.",
            )
            ph_test_results["Alpha"] = config.ph_test_alpha
        except Exception as exc:
            ph_test_results = pd.DataFrame(
                [{
                    "Feature": "ALL",
                    "Test_Statistic": np.nan,
                    "P_Value": np.nan,
                    "P_Value_Holm": np.nan,
                    "P_Value_BH": np.nan,
                    "Raw_PH_Status": "TEST_ERROR",
                    "Holm_PH_Status": "TEST_ERROR",
                    "BH_PH_Status": "TEST_ERROR",
                    "PH_Assumption_Status": "TEST_ERROR",
                    "Review_Rationale": str(exc),
                    "Alpha": config.ph_test_alpha,
                }]
            )
    else:
        ph_test_results = pd.DataFrame(
            [{
                "Feature": "ALL",
                "Test_Statistic": np.nan,
                "P_Value": np.nan,
                "P_Value_Holm": np.nan,
                "P_Value_BH": np.nan,
                "Raw_PH_Status": "NOT_RUN",
                "Holm_PH_Status": "NOT_RUN",
                "BH_PH_Status": "NOT_RUN",
                "PH_Assumption_Status": "NOT_RUN",
                "Review_Rationale": "PH testing was disabled by configuration.",
                "Alpha": config.ph_test_alpha,
            }]
        )

    if ph_test_results["PH_Assumption_Status"].eq("TEST_ERROR").any():
        ph_global_status = "REVIEW"
    elif ph_test_results["Holm_PH_Status"].eq("REVIEW").any():
        ph_global_status = "REVIEW"
    elif ph_test_results["Raw_PH_Status"].eq("REVIEW").any():
        ph_global_status = "PASS_WITH_REVIEW"
    elif ph_test_results["PH_Assumption_Status"].eq("NOT_RUN").any():
        ph_global_status = "NOT_RUN"
    else:
        ph_global_status = "PASS"

    review_features = ph_test_results.loc[
        ph_test_results["Raw_PH_Status"].eq("REVIEW"), "Feature"
    ].astype(str).tolist()
    ph_residuals = _extract_scaled_schoenfeld_residuals(
        cph,
        training_frame,
        config.duration_col,
        review_features,
    )

    cv_pass = cv_results["Status"].eq("PASS")
    cv_values = pd.to_numeric(
        cv_results.loc[cv_pass, "Concordance_Index"], errors="coerce"
    ).dropna()
    calibration_valid = calibration_metrics.get(
        "Calibration_Status", pd.Series(dtype=str)
    ).eq("PASS").all()

    fit_metadata = pd.DataFrame(
        [{
            "Engine_Release": ENGINE_RELEASE,
            "Rows": int(len(training_frame)),
            "Events": int(training_frame[config.event_col].sum()),
            "Censored": int(len(training_frame) - training_frame[config.event_col].sum()),
            "Event_Rate": float(training_frame[config.event_col].mean()),
            "Feature_Count": int(len(config.model_features)),
            "Penalizer": float(config.penalizer),
            "L1_Ratio": float(config.l1_ratio),
            "Standardized_Features": bool(config.standardize_model_features),
            "Apparent_Concordance_Index": float(cph.concordance_index_),
            "CV_Folds": int(config.cross_validation_folds),
            "CV_Mean_Concordance": float(cv_values.mean()) if not cv_values.empty else np.nan,
            "CV_SD_Concordance": float(cv_values.std(ddof=1)) if len(cv_values) > 1 else 0.0,
            "CV_Min_Concordance": float(cv_values.min()) if not cv_values.empty else np.nan,
            "CV_Max_Concordance": float(cv_values.max()) if not cv_values.empty else np.nan,
            "CV_Status": "PASS" if cv_pass.all() and len(cv_values) == config.cross_validation_folds else "REVIEW",
            "Integrated_Brier_Score": (
                float(calibration_metrics["Integrated_Brier_Score"].iloc[0])
                if "Integrated_Brier_Score" in calibration_metrics.columns
                else np.nan
            ),
            "Max_Group_Absolute_Calibration_Error": (
                float(calibration_metrics["Max_Group_Absolute_Calibration_Error"].iloc[0])
                if "Max_Group_Absolute_Calibration_Error" in calibration_metrics.columns
                else np.nan
            ),
            "Calibration_Status": "PASS" if calibration_valid else "REVIEW",
            "PH_Global_Status": ph_global_status,
            "Partial_AIC": float(cph.AIC_partial_),
            "Log_Likelihood": float(cph.log_likelihood_),
            "Fit_Status": "SUCCESS",
        }]
    )

    fit_warnings = pd.DataFrame(
        [{"Warning": message} for message in captured_warnings]
        or [{"Warning": "No fit warnings were captured."}]
    )

    return CoxModelArtifacts(
        model=cph,
        model_scaler=model_scaler,
        training_frame=training_frame,
        model_summary=model_summary,
        fit_metadata=fit_metadata,
        fit_warnings=fit_warnings,
        cross_validation_results=cv_results,
        cross_validated_predictions=cv_predictions,
        calibration_metrics=calibration_metrics,
        calibration_by_risk_group=calibration_by_group,
        ph_test_results=ph_test_results,
        ph_residuals=ph_residuals,
        ph_sensitivity_results=pd.DataFrame(),
    )

# ==================================================================================================
# POPULATION RISK SCORING AND GOVERNED TARGET-COHORT SELECTION
# ==================================================================================================

def score_and_stratify_population(
    df: pd.DataFrame,
    personas: PersonaArtifacts,
    model_artifacts: CoxModelArtifacts,
    config: FrameworkConfig,
) -> RiskArtifacts:
    """Scores the full population and selects the governed baseline target cohort."""
    model_matrix = transform_model_features(
        df,
        config,
        model_artifacts.model_scaler,
    )
    partial_hazard = model_artifacts.model.predict_partial_hazard(model_matrix)
    hazard_series = pd.Series(
        np.asarray(partial_hazard).reshape(-1),
        index=df.index,
        name="Partial_Hazard",
    )

    scored = df.copy()
    scored["Persona"] = personas.scored_data["Persona"]
    scored["Segment_ID"] = personas.scored_data["Segment_ID"]
    scored["Partial_Hazard"] = hazard_series
    scored["Risk_Percentile"] = scored["Partial_Hazard"].rank(
        method="average",
        pct=True,
    )

    q25 = float(scored["Partial_Hazard"].quantile(0.25))
    target_threshold = float(scored["Partial_Hazard"].quantile(config.risk_quantile))
    target_label = f"High Risk (Top {1.0 - config.risk_quantile:.0%})"

    scored["Risk_Tier"] = "Mid Risk"
    scored.loc[scored["Partial_Hazard"] <= q25, "Risk_Tier"] = "Benchmark (Bottom 25%)"
    scored.loc[
        scored["Partial_Hazard"] >= target_threshold,
        "Risk_Tier",
    ] = target_label
    scored["Target_Cohort_Flag"] = (
        scored["Partial_Hazard"] >= target_threshold
    ).astype(int)

    profile_features = config.profile_features or config.segmentation_features
    grouped = scored.groupby("Risk_Tier", sort=False)
    risk_profiles = grouped.agg(
        Record_Count=(config.id_col, "count"),
        Event_Rate=(config.event_col, "mean"),
        Median_Duration=(config.duration_col, "median"),
        Mean_Partial_Hazard=("Partial_Hazard", "mean"),
    ).reset_index()
    feature_means = grouped[list(profile_features)].mean().reset_index()
    risk_profiles = risk_profiles.merge(feature_means, on="Risk_Tier", how="left")

    target_index = scored.index[scored["Target_Cohort_Flag"] == 1]
    target_columns = [
        config.id_col,
        "Segment_ID",
        "Persona",
        "Risk_Tier",
        "Partial_Hazard",
        "Risk_Percentile",
        *[feature for feature in profile_features if feature in scored.columns],
    ]
    target_cohort = scored.loc[target_index, list(dict.fromkeys(target_columns))].copy()

    return RiskArtifacts(
        scored_data=scored,
        risk_tier_profiles=risk_profiles,
        target_cohort=target_cohort,
        target_index=target_index,
        target_threshold=target_threshold,
        target_label=target_label,
    )


def run_ph_sensitivity_analysis(
    source_df: pd.DataFrame,
    config: FrameworkConfig,
    model_artifacts: CoxModelArtifacts,
    risk_artifacts: RiskArtifacts,
    scenarios_config: Sequence[ScenarioDefinition] = (),
    primary_scenario_summary: pd.DataFrame | None = None,
) -> pd.DataFrame:
    """Fits stratified sensitivity models for low-cardinality raw PH review features.

    This does not automatically replace the primary model. It answers a narrower validation
    question: if a flagged low-cardinality feature is moved from the linear predictor into a strata
    definition, do discrimination, the selected target cohort, and the ranking of favorable scenarios
    remain materially stable?
    """
    raw_review = model_artifacts.ph_test_results.loc[
        model_artifacts.ph_test_results["Raw_PH_Status"].eq("REVIEW")
    ]
    if raw_review.empty:
        return pd.DataFrame(
            [{
                "Feature": "NONE",
                "Sensitivity_Method": "NOT_APPLICABLE",
                "Fit_Status": "PASS",
                "Conclusion": "No raw PH feature review items required sensitivity fitting.",
            }]
        )

    transformed_full = transform_model_features(
        source_df,
        config,
        model_artifacts.model_scaler,
    )
    baseline_target_ids = set(risk_artifacts.target_cohort[config.id_col].astype(str))
    baseline_cindex = float(
        model_artifacts.fit_metadata.iloc[0]["Apparent_Concordance_Index"]
    )
    primary_summary = (
        primary_scenario_summary.copy()
        if primary_scenario_summary is not None
        else pd.DataFrame()
    )
    primary_improvement_order = (
        primary_summary.loc[
            primary_summary["Scenario_Type"].eq("IMPROVEMENT")
        ]
        .sort_values("Improvement_Rank")["Scenario_Name"]
        .astype(str)
        .tolist()
        if not primary_summary.empty and "Improvement_Rank" in primary_summary.columns
        else []
    )
    rows: list[dict[str, Any]] = []

    for _, ph_row in raw_review.iterrows():
        feature = str(ph_row["Feature"])
        unique_levels = int(source_df[feature].nunique(dropna=True))
        base = {
            "Feature": feature,
            "Raw_P_Value": float(ph_row["P_Value"]),
            "Holm_Adjusted_P_Value": float(ph_row["P_Value_Holm"]),
            "Holm_PH_Status": str(ph_row["Holm_PH_Status"]),
            "Feature_Unique_Levels": unique_levels,
        }

        if unique_levels > config.ph_strata_max_unique:
            rows.append(
                {
                    **base,
                    "Sensitivity_Method": "MANUAL_TIME_INTERACTION_REVIEW",
                    "Fit_Status": "NOT_RUN",
                    "Apparent_Concordance_Index": np.nan,
                    "Concordance_Delta": np.nan,
                    "Target_Cohort_Overlap_Rate": np.nan,
                    "Target_Cohort_Jaccard": np.nan,
                    "Scenario_Rank_Spearman": np.nan,
                    "Scenario_Rank_Exact_Match": np.nan,
                    "Primary_Best_Improvement": (
                        primary_improvement_order[0] if primary_improvement_order else ""
                    ),
                    "Sensitivity_Best_Improvement": "",
                    "Scenario_Rank_Status": "NOT_RUN",
                    "Conclusion": (
                        "Feature has too many levels for automatic stratification; inspect scaled "
                        "Schoenfeld residuals and consider a governed time interaction."
                    ),
                }
            )
            continue

        try:
            strata_column = f"__strata__{feature}"
            sensitivity_features = [
                item for item in config.model_features if item != feature
            ]
            training = transformed_full.loc[:, sensitivity_features].copy()
            training[strata_column] = source_df[feature].astype(str)
            training[config.duration_col] = source_df[config.duration_col]
            training[config.event_col] = source_df[config.event_col].astype(int)

            sensitivity_model = CoxPHFitter(
                penalizer=config.penalizer,
                l1_ratio=config.l1_ratio,
            )
            sensitivity_model.fit(
                training,
                duration_col=config.duration_col,
                event_col=config.event_col,
                strata=[strata_column],
                show_progress=False,
            )

            prediction_frame = transformed_full.loc[:, sensitivity_features].copy()
            prediction_frame[strata_column] = source_df[feature].astype(str)
            sensitivity_hazard = np.asarray(
                sensitivity_model.predict_partial_hazard(prediction_frame)
            ).reshape(-1)
            threshold = float(np.quantile(sensitivity_hazard, config.risk_quantile))
            sensitivity_ids = set(
                source_df.loc[sensitivity_hazard >= threshold, config.id_col].astype(str)
            )
            intersection = len(baseline_target_ids.intersection(sensitivity_ids))
            union = len(baseline_target_ids.union(sensitivity_ids))

            # Re-score the primary governed target IDs under the stratified sensitivity model. The
            # purpose is not to replace the primary model automatically; it is to test whether the same
            # business scenario ranking survives a reasonable alternative PH treatment.
            scenario_reductions: dict[str, float] = {}
            if scenarios_config:
                sensitivity_baseline = prediction_frame.loc[risk_artifacts.target_index]
                sensitivity_baseline_hazard = np.asarray(
                    sensitivity_model.predict_partial_hazard(sensitivity_baseline)
                ).reshape(-1)
                sensitivity_baseline_mean = float(sensitivity_baseline_hazard.mean())

                for scenario in scenarios_config:
                    if scenario.scenario_type != "IMPROVEMENT":
                        continue
                    scenario_df = source_df.copy()
                    for change in scenario.changes:
                        scenario_df, _ = apply_feature_change(
                            scenario_df,
                            risk_artifacts.target_index,
                            change,
                            scenario,
                        )
                    scenario_df, _ = sync_engineered_dependencies(
                        scenario_df,
                        config.model_features,
                        scenario.name,
                    )
                    transformed_scenario = transform_model_features(
                        scenario_df,
                        config,
                        model_artifacts.model_scaler,
                    )
                    sensitivity_scenario = transformed_scenario.loc[
                        :, sensitivity_features
                    ].copy()
                    sensitivity_scenario[strata_column] = scenario_df[feature].astype(str)
                    scenario_hazard = np.asarray(
                        sensitivity_model.predict_partial_hazard(
                            sensitivity_scenario.loc[risk_artifacts.target_index]
                        )
                    ).reshape(-1)
                    scenario_reductions[scenario.name] = (
                        (sensitivity_baseline_mean - float(scenario_hazard.mean()))
                        / sensitivity_baseline_mean
                        if sensitivity_baseline_mean != 0
                        else np.nan
                    )

            sensitivity_order = [
                name
                for name, _ in sorted(
                    scenario_reductions.items(),
                    key=lambda item: item[1],
                    reverse=True,
                )
            ]
            common_names = [
                name for name in primary_improvement_order if name in scenario_reductions
            ]
            if len(common_names) >= 2:
                primary_ranks = pd.Series(
                    range(1, len(common_names) + 1),
                    index=common_names,
                    dtype=float,
                )
                sensitivity_ranks = pd.Series(
                    {
                        name: rank
                        for rank, name in enumerate(sensitivity_order, start=1)
                        if name in common_names
                    },
                    dtype=float,
                ).reindex(common_names)
                rank_spearman = float(primary_ranks.corr(sensitivity_ranks, method="spearman"))
                exact_match = sensitivity_order[: len(common_names)] == common_names
            elif len(common_names) == 1:
                rank_spearman = 1.0
                exact_match = sensitivity_order[:1] == common_names
            else:
                rank_spearman = np.nan
                exact_match = np.nan

            ranking_status = (
                "PASS"
                if (
                    not primary_improvement_order
                    or (
                        np.isfinite(rank_spearman)
                        and rank_spearman >= 0.90
                        and bool(exact_match)
                    )
                )
                else "REVIEW"
            )

            rows.append(
                {
                    **base,
                    "Sensitivity_Method": "STRATIFIED_COXPH",
                    "Fit_Status": "PASS",
                    "Apparent_Concordance_Index": float(
                        sensitivity_model.concordance_index_
                    ),
                    "Concordance_Delta": float(
                        sensitivity_model.concordance_index_ - baseline_cindex
                    ),
                    "Target_Cohort_Overlap_Rate": (
                        intersection / len(baseline_target_ids)
                        if baseline_target_ids
                        else np.nan
                    ),
                    "Target_Cohort_Jaccard": intersection / union if union else np.nan,
                    "Scenario_Rank_Spearman": rank_spearman,
                    "Scenario_Rank_Exact_Match": exact_match,
                    "Primary_Best_Improvement": (
                        primary_improvement_order[0] if primary_improvement_order else ""
                    ),
                    "Sensitivity_Best_Improvement": (
                        sensitivity_order[0] if sensitivity_order else ""
                    ),
                    "Scenario_Rank_Status": ranking_status,
                    "Conclusion": (
                        "Automatic stratified sensitivity fit completed. Review discrimination, "
                        "target-cohort overlap, and improvement-scenario ranking before changing "
                        "the primary specification."
                    ),
                }
            )
        except Exception as exc:
            rows.append(
                {
                    **base,
                    "Sensitivity_Method": "STRATIFIED_COXPH",
                    "Fit_Status": "FAIL",
                    "Apparent_Concordance_Index": np.nan,
                    "Concordance_Delta": np.nan,
                    "Target_Cohort_Overlap_Rate": np.nan,
                    "Target_Cohort_Jaccard": np.nan,
                    "Scenario_Rank_Spearman": np.nan,
                    "Scenario_Rank_Exact_Match": np.nan,
                    "Primary_Best_Improvement": (
                        primary_improvement_order[0] if primary_improvement_order else ""
                    ),
                    "Sensitivity_Best_Improvement": "",
                    "Scenario_Rank_Status": "NOT_AVAILABLE",
                    "Conclusion": str(exc),
                }
            )

    return pd.DataFrame(rows)

# ==================================================================================================
# ENGINEERED-FEATURE DEPENDENCY SYNCHRONIZATION
# ==================================================================================================

def sync_engineered_dependencies(
    df: pd.DataFrame,
    model_features: Sequence[str],
    scenario_name: str,
) -> tuple[pd.DataFrame, pd.DataFrame]:
    """Rebuilds interaction and squared terms after a scenario changes base features.

    The audit records both the pre-rebuild discrepancy and the post-rebuild verification. This
    turns dependency synchronization from an undocumented convenience into explicit control
    evidence.
    """
    synced = df.copy()
    audit_rows: list[dict[str, Any]] = []

    for feature in model_features:
        sources = engineered_feature_sources(feature)
        if not sources:
            continue

        if any(source not in synced.columns for source in sources):
            audit_rows.append(
                {
                    "Scenario_Name": scenario_name,
                    "Engineered_Feature": feature,
                    "Dependency_Type": "INTERACTION" if "_x_" in feature else "SQUARE",
                    "Source_Features": " | ".join(sources),
                    "Pre_Rebuild_Max_Abs_Diff": np.nan,
                    "Post_Rebuild_Max_Abs_Diff": np.nan,
                    "Recalculated": False,
                    "Validation_Status": "FAIL_MISSING_SOURCE",
                }
            )
            continue

        previous = pd.to_numeric(synced[feature], errors="coerce")
        if "_x_" in feature:
            expected = synced.loc[:, list(sources)].prod(axis=1)
            dependency_type = "INTERACTION"
        else:
            expected = pd.to_numeric(synced[sources[0]], errors="coerce") ** 2
            dependency_type = "SQUARE"

        pre_difference = float(np.nanmax(np.abs(previous.to_numpy() - expected.to_numpy())))
        synced[feature] = expected
        post_difference = float(
            np.nanmax(
                np.abs(
                    pd.to_numeric(synced[feature], errors="coerce").to_numpy()
                    - expected.to_numpy()
                )
            )
        )

        audit_rows.append(
            {
                "Scenario_Name": scenario_name,
                "Engineered_Feature": feature,
                "Dependency_Type": dependency_type,
                "Source_Features": " | ".join(sources),
                "Pre_Rebuild_Max_Abs_Diff": pre_difference,
                "Post_Rebuild_Max_Abs_Diff": post_difference,
                "Recalculated": True,
                "Validation_Status": "PASS" if post_difference <= 1e-10 else "FAIL",
            }
        )

    audit = pd.DataFrame(audit_rows)
    if audit.empty:
        audit = pd.DataFrame(
            [
                {
                    "Scenario_Name": scenario_name,
                    "Engineered_Feature": "NONE",
                    "Dependency_Type": "NONE",
                    "Source_Features": "",
                    "Pre_Rebuild_Max_Abs_Diff": 0.0,
                    "Post_Rebuild_Max_Abs_Diff": 0.0,
                    "Recalculated": False,
                    "Validation_Status": "NOT_APPLICABLE",
                }
            ]
        )
    return synced, audit


# ==================================================================================================
# GENERIC SAME-COHORT SCENARIO SIMULATION
# ==================================================================================================

def apply_feature_change(
    df: pd.DataFrame,
    row_index: pd.Index,
    change: FeatureChange,
    scenario: ScenarioDefinition,
) -> tuple[pd.DataFrame, dict[str, Any]]:
    """Applies one bounded change and returns explicit configured-versus-realized evidence."""
    result = df.copy()
    baseline = pd.to_numeric(
        result.loc[row_index, change.feature], errors="raise"
    ).astype(float)

    if change.operation == "multiply":
        unconstrained = baseline * change.value
    elif change.operation == "add":
        unconstrained = baseline + change.value
    elif change.operation == "replace":
        unconstrained = pd.Series(change.value, index=baseline.index, dtype=float)
    else:
        raise ValueError(f"Unsupported scenario operation: {change.operation}")

    changed = unconstrained.copy()
    if change.lower_bound is not None:
        changed = changed.clip(lower=change.lower_bound)
    if change.upper_bound is not None:
        changed = changed.clip(upper=change.upper_bound)
    if change.round_digits is not None:
        changed = changed.round(change.round_digits)

    result.loc[row_index, change.feature] = changed
    difference = changed - baseline
    nonzero_denominator = baseline.abs() > 1e-12
    relative_change = pd.Series(np.nan, index=baseline.index, dtype=float)
    relative_change.loc[nonzero_denominator] = (
        difference.loc[nonzero_denominator]
        / baseline.loc[nonzero_denominator].abs()
    )

    target_count = int(len(row_index))
    records_changed = int((difference.abs() > 1e-12).sum())
    lower_bound_hits = int(
        np.isclose(changed, change.lower_bound).sum()
        if change.lower_bound is not None
        else 0
    )
    upper_bound_hits = int(
        np.isclose(changed, change.upper_bound).sum()
        if change.upper_bound is not None
        else 0
    )
    audit = {
        "Scenario_Name": scenario.name,
        "Scenario_Type": scenario.scenario_type,
        "Expected_Direction": scenario.expected_direction,
        "Feature": change.feature,
        "Operation": change.operation,
        "Configured_Value": change.value,
        "Lower_Bound": change.lower_bound,
        "Upper_Bound": change.upper_bound,
        "Round_Digits": change.round_digits,
        "Target_Cohort_Count": target_count,
        "Records_Changed": records_changed,
        "Records_Changed_Rate": records_changed / target_count if target_count else np.nan,
        "Records_Hitting_Lower_Bound": lower_bound_hits,
        "Lower_Bound_Hit_Rate": lower_bound_hits / target_count if target_count else np.nan,
        "Records_Hitting_Upper_Bound": upper_bound_hits,
        "Upper_Bound_Hit_Rate": upper_bound_hits / target_count if target_count else np.nan,
        "Mean_Baseline_Value": float(baseline.mean()),
        "Mean_Unconstrained_Value": float(unconstrained.mean()),
        "Mean_Scenario_Value": float(changed.mean()),
        "Mean_Absolute_Change": float(difference.abs().mean()),
        "Mean_Signed_Change": float(difference.mean()),
        "Mean_Relative_Change": float(relative_change.mean(skipna=True)),
        "Configured_vs_Realized_Note": (
            "Bounds and rounding may cause realized movement to differ from the configured change."
        ),
    }
    return result, audit


def mean_predicted_survival(
    model: CoxPHFitter,
    model_matrix: pd.DataFrame,
    times: np.ndarray,
) -> pd.Series:
    """Returns cohort-average CoxPH survival and enforces the mathematical S(0)=1 identity."""
    requested_times = np.asarray(times, dtype=float)
    survival = model.predict_survival_function(model_matrix, times=requested_times)
    mean_curve = pd.Series(
        survival.mean(axis=1).to_numpy(dtype=float),
        index=requested_times,
        name="Mean_Predicted_Survival",
    )
    zero_mask = np.isclose(mean_curve.index.to_numpy(dtype=float), 0.0)
    if zero_mask.any():
        mean_curve.iloc[np.where(zero_mask)[0]] = 1.0
    return mean_curve


def simulate_scenarios(
    source_df: pd.DataFrame,
    config: FrameworkConfig,
    scenarios: Sequence[ScenarioDefinition],
    model_artifacts: CoxModelArtifacts,
    risk_artifacts: RiskArtifacts,
) -> ScenarioArtifacts:
    """Applies controlled scenarios and re-scores the identical baseline target cohort."""
    target_index = risk_artifacts.target_index
    if len(target_index) == 0:
        raise ValueError("Target cohort is empty; revise risk_quantile or input data.")

    baseline_matrix_all = transform_model_features(
        source_df,
        config,
        model_artifacts.model_scaler,
    )
    baseline_matrix = baseline_matrix_all.loc[target_index]

    baseline_hazard = pd.Series(
        np.asarray(model_artifacts.model.predict_partial_hazard(baseline_matrix)).reshape(-1),
        index=target_index,
        name="Baseline_Partial_Hazard",
    )
    baseline_mean_hazard = float(baseline_hazard.mean())

    max_time = max(float(source_df[config.duration_col].max()), config.evaluation_horizon)
    timeline = np.linspace(0.0, max_time, config.timeline_points)
    baseline_curve = mean_predicted_survival(
        model_artifacts.model,
        baseline_matrix,
        timeline,
    )
    baseline_horizon = float(
        mean_predicted_survival(
            model_artifacts.model,
            baseline_matrix,
            np.array([config.evaluation_horizon]),
        ).iloc[0]
    )

    summary_rows: list[dict[str, Any]] = []
    score_frames: list[pd.DataFrame] = []
    change_audits: list[dict[str, Any]] = []
    dependency_frames: list[pd.DataFrame] = []
    curve_frame = pd.DataFrame(
        {
            "Time": timeline,
            "Baseline Target Cohort": baseline_curve.to_numpy(),
        }
    )

    for scenario in scenarios:
        scenario_df = source_df.copy()
        if not scenario.changes:
            change_audits.append(
                {
                    "Scenario_Name": scenario.name,
                    "Scenario_Type": scenario.scenario_type,
                    "Expected_Direction": scenario.expected_direction,
                    "Feature": "NO_CHANGE",
                    "Operation": "NONE",
                    "Configured_Value": 0.0,
                    "Lower_Bound": np.nan,
                    "Upper_Bound": np.nan,
                    "Round_Digits": np.nan,
                    "Target_Cohort_Count": int(len(target_index)),
                    "Records_Changed": 0,
                    "Records_Changed_Rate": 0.0,
                    "Records_Hitting_Lower_Bound": 0,
                    "Lower_Bound_Hit_Rate": 0.0,
                    "Records_Hitting_Upper_Bound": 0,
                    "Upper_Bound_Hit_Rate": 0.0,
                    "Mean_Baseline_Value": np.nan,
                    "Mean_Unconstrained_Value": np.nan,
                    "Mean_Scenario_Value": np.nan,
                    "Mean_Absolute_Change": 0.0,
                    "Mean_Signed_Change": 0.0,
                    "Mean_Relative_Change": 0.0,
                    "Configured_vs_Realized_Note": "No-change control; no feature movement configured.",
                }
            )
        for change in scenario.changes:
            scenario_df, change_audit = apply_feature_change(
                scenario_df,
                target_index,
                change,
                scenario,
            )
            change_audits.append(change_audit)

        scenario_df, dependency_audit = sync_engineered_dependencies(
            scenario_df,
            config.model_features,
            scenario.name,
        )
        dependency_frames.append(dependency_audit)

        scenario_validation = validate_survival_input(scenario_df, config, scenarios=())
        scenario_validation.raise_for_errors()

        scenario_matrix_all = transform_model_features(
            scenario_df,
            config,
            model_artifacts.model_scaler,
        )
        scenario_matrix = scenario_matrix_all.loc[target_index]

        scenario_hazard = pd.Series(
            np.asarray(model_artifacts.model.predict_partial_hazard(scenario_matrix)).reshape(-1),
            index=target_index,
            name="Scenario_Partial_Hazard",
        )
        scenario_mean_hazard = float(scenario_hazard.mean())
        relative_hazard_reduction = (
            (baseline_mean_hazard - scenario_mean_hazard) / baseline_mean_hazard
            if baseline_mean_hazard != 0
            else np.nan
        )

        scenario_curve = mean_predicted_survival(
            model_artifacts.model,
            scenario_matrix,
            timeline,
        )
        curve_frame[scenario.name] = scenario_curve.to_numpy()

        scenario_horizon = float(
            mean_predicted_survival(
                model_artifacts.model,
                scenario_matrix,
                np.array([config.evaluation_horizon]),
            ).iloc[0]
        )
        survival_uplift = scenario_horizon - baseline_horizon
        direction = modeled_direction(
            relative_hazard_reduction,
            survival_uplift,
            config.neutral_scenario_tolerance,
        )
        direction_status = (
            "PASS"
            if scenario.expected_direction in {"UNSPECIFIED", direction}
            else "FAIL"
        )

        summary_rows.append(
            {
                "Scenario_Name": scenario.name,
                "Scenario_Type": scenario.scenario_type,
                "Expected_Direction": scenario.expected_direction,
                "Description": scenario.description,
                "Configured_Change_Count": int(len(scenario.changes)),
                "Target_Cohort_Count": int(len(target_index)),
                "Baseline_Mean_Partial_Hazard": baseline_mean_hazard,
                "Scenario_Mean_Partial_Hazard": scenario_mean_hazard,
                "Relative_Hazard_Reduction": relative_hazard_reduction,
                "Evaluation_Horizon": float(config.evaluation_horizon),
                "Time_Unit": config.time_unit_label,
                "Baseline_Survival_At_Horizon": baseline_horizon,
                "Scenario_Survival_At_Horizon": scenario_horizon,
                "Survival_Probability_Uplift": survival_uplift,
                "Survival_Uplift_Percentage_Points": survival_uplift * 100.0,
                "Modeled_Direction": direction,
                "Direction_Check_Status": direction_status,
            }
        )

        score_frame = pd.DataFrame(
            {
                config.id_col: source_df.loc[target_index, config.id_col].to_numpy(),
                "Scenario_Name": scenario.name,
                "Scenario_Type": scenario.scenario_type,
                "Baseline_Partial_Hazard": baseline_hazard.to_numpy(),
                "Scenario_Partial_Hazard": scenario_hazard.to_numpy(),
                "Partial_Hazard_Delta": (
                    scenario_hazard.to_numpy() - baseline_hazard.to_numpy()
                ),
            },
            index=target_index,
        )
        score_frame["Relative_Partial_Hazard_Change"] = np.where(
            score_frame["Baseline_Partial_Hazard"] != 0,
            score_frame["Partial_Hazard_Delta"]
            / score_frame["Baseline_Partial_Hazard"],
            np.nan,
        )
        score_frames.append(score_frame.reset_index(drop=True))

    summary = pd.DataFrame(summary_rows)
    if not summary.empty:
        summary = summary.sort_values(
            ["Relative_Hazard_Reduction", "Survival_Probability_Uplift"],
            ascending=[False, False],
        ).reset_index(drop=True)
        summary["Scenario_Rank"] = np.arange(1, len(summary) + 1)
        summary["Improvement_Rank"] = np.nan
        improvement_mask = summary["Scenario_Type"].eq("IMPROVEMENT")
        summary.loc[improvement_mask, "Improvement_Rank"] = np.arange(
            1,
            int(improvement_mask.sum()) + 1,
        )

    return ScenarioArtifacts(
        summary=summary,
        target_scores=(
            pd.concat(score_frames, ignore_index=True)
            if score_frames
            else pd.DataFrame()
        ),
        scenario_change_audit=pd.DataFrame(change_audits),
        dependency_audit=(
            pd.concat(dependency_frames, ignore_index=True)
            if dependency_frames
            else pd.DataFrame()
        ),
        predicted_survival_curves=curve_frame,
    )

# ==================================================================================================
# EVIDENCE VISUALS
# ==================================================================================================

def _prepare_dark_axis(ax: plt.Axes) -> None:
    ax.set_facecolor(COLOR_BACKGROUND)
    ax.tick_params(colors=COLOR_LIGHT_GRAY)
    for spine in ax.spines.values():
        spine.set_color(COLOR_MID_GRAY)
    ax.xaxis.label.set_color(COLOR_LIGHT_GRAY)
    ax.yaxis.label.set_color(COLOR_LIGHT_GRAY)
    ax.title.set_color(COLOR_WHITE)
    ax.grid(alpha=0.15, color=COLOR_MID_GRAY)


def plot_persona_kaplan_meier(
    data: pd.DataFrame,
    config: FrameworkConfig,
    output_path: Path,
) -> None:
    figure, axis = plt.subplots(figsize=(10, 6), facecolor=COLOR_BACKGROUND)
    _prepare_dark_axis(axis)
    kmf = KaplanMeierFitter()

    for persona in sorted(data["Persona"].unique()):
        subset = data[data["Persona"] == persona]
        kmf.fit(
            subset[config.duration_col],
            event_observed=subset[config.event_col],
            label=persona,
        )
        kmf.plot_survival_function(
            ax=axis,
            ci_show=config.show_km_confidence_intervals,
        )

    axis.set_title("Persona Kaplan-Meier Survival Evidence")
    axis.set_xlabel(f"Time ({config.time_unit_label})")
    axis.set_ylabel("Observed Probability of Remaining Event-Free")
    legend = axis.legend(facecolor=COLOR_PANEL, edgecolor=COLOR_MID_GRAY)
    for item in legend.get_texts():
        item.set_color(COLOR_WHITE)
    figure.tight_layout()
    figure.savefig(output_path, dpi=180, bbox_inches="tight", facecolor=figure.get_facecolor())
    plt.close(figure)


def plot_risk_tier_kaplan_meier(
    data: pd.DataFrame,
    config: FrameworkConfig,
    output_path: Path,
) -> None:
    figure, axis = plt.subplots(figsize=(10, 6), facecolor=COLOR_BACKGROUND)
    _prepare_dark_axis(axis)
    kmf = KaplanMeierFitter()

    tier_order = [
        f"High Risk (Top {1.0 - config.risk_quantile:.0%})",
        "Mid Risk",
        "Benchmark (Bottom 25%)",
    ]
    for tier in tier_order:
        subset = data[data["Risk_Tier"] == tier]
        if subset.empty:
            continue
        kmf.fit(
            subset[config.duration_col],
            event_observed=subset[config.event_col],
            label=tier,
        )
        kmf.plot_survival_function(
            ax=axis,
            ci_show=config.show_km_confidence_intervals,
        )

    axis.set_title("Observed Kaplan-Meier Survival by Model-Derived Risk Tier")
    axis.set_xlabel(f"Time ({config.time_unit_label})")
    axis.set_ylabel("Observed Probability of Remaining Event-Free")
    legend = axis.legend(facecolor=COLOR_PANEL, edgecolor=COLOR_MID_GRAY)
    for item in legend.get_texts():
        item.set_color(COLOR_WHITE)
    figure.tight_layout()
    figure.savefig(output_path, dpi=180, bbox_inches="tight", facecolor=figure.get_facecolor())
    plt.close(figure)


def scenario_color(scenario_type: str, ordinal: int) -> str:
    """Returns a consistent semantic color for scenario plots."""
    if scenario_type == "CONTROL":
        return COLOR_MID_GRAY
    if scenario_type == "STRESS":
        return COLOR_RED
    if scenario_type == "TECHNICAL":
        return COLOR_GOLD
    return [COLOR_CYAN, COLOR_GOLD, COLOR_TEAL, COLOR_GREEN][ordinal % 4]


def plot_predicted_survival_comparison(
    curves: pd.DataFrame,
    scenario_summary: pd.DataFrame,
    config: FrameworkConfig,
    output_path: Path,
) -> None:
    figure, axis = plt.subplots(figsize=(10, 6), facecolor=COLOR_BACKGROUND)
    _prepare_dark_axis(axis)

    axis.plot(
        curves["Time"],
        curves["Baseline Target Cohort"],
        label="Baseline Target Cohort",
        linewidth=3.2,
        color=COLOR_LIGHT_GRAY,
    )

    ordered = scenario_summary.sort_values("Scenario_Rank") if not scenario_summary.empty else pd.DataFrame()
    for ordinal, (_, row) in enumerate(ordered.iterrows()):
        scenario_name = str(row["Scenario_Name"])
        if scenario_name not in curves.columns:
            continue
        axis.plot(
            curves["Time"],
            curves[scenario_name],
            label=scenario_name,
            linewidth=2.1,
            color=scenario_color(str(row["Scenario_Type"]), ordinal),
        )

    axis.set_title("CoxPH-Predicted Survival: Same Target Cohort, Controlled Scenarios")
    axis.set_xlabel(f"Time ({config.time_unit_label})")
    axis.set_ylabel("Mean Predicted Survival Probability")
    axis.set_ylim(0.0, 1.02)
    legend = axis.legend(facecolor=COLOR_PANEL, edgecolor=COLOR_MID_GRAY, fontsize=8)
    for item in legend.get_texts():
        item.set_color(COLOR_WHITE)
    figure.tight_layout()
    figure.savefig(output_path, dpi=180, bbox_inches="tight", facecolor=figure.get_facecolor())
    plt.close(figure)


def plot_scenario_hazard_reduction(
    summary: pd.DataFrame,
    output_path: Path,
) -> None:
    figure, axis = plt.subplots(figsize=(10.5, 6.2), facecolor=COLOR_BACKGROUND)
    _prepare_dark_axis(axis)

    if summary.empty:
        axis.text(0.5, 0.5, "No configured scenarios", ha="center", va="center", color=COLOR_WHITE)
    else:
        ordered = summary.sort_values("Relative_Hazard_Reduction", ascending=True)
        values = ordered["Relative_Hazard_Reduction"].to_numpy(dtype=float) * 100.0
        colors_for_bars = [
            scenario_color(str(scenario_type), index)
            for index, scenario_type in enumerate(ordered["Scenario_Type"])
        ]
        axis.barh(ordered["Scenario_Name"], values, color=colors_for_bars, alpha=0.88)
        axis.axvline(0.0, color=COLOR_LIGHT_GRAY, linewidth=1.1)

        span = max(1.0, float(np.nanmax(np.abs(values))))
        for position, value in enumerate(values):
            axis.text(
                value + (0.012 * span if value >= 0 else -0.012 * span),
                position,
                f"{value:.1f}%",
                va="center",
                ha="left" if value >= 0 else "right",
                color=COLOR_WHITE,
                fontsize=9,
            )

    axis.set_title("Modeled Relative Hazard Movement by Scenario")
    axis.set_xlabel("Relative Reduction in Mean Partial Hazard (%) — negative values are adverse")
    figure.subplots_adjust(left=0.31, right=0.97, top=0.90, bottom=0.14)
    figure.savefig(output_path, dpi=180, bbox_inches="tight", facecolor=figure.get_facecolor())
    plt.close(figure)


def plot_calibration_evidence(
    calibration_by_group: pd.DataFrame,
    config: FrameworkConfig,
    output_path: Path,
) -> None:
    figure, axis = plt.subplots(figsize=(7.8, 6.4), facecolor=COLOR_BACKGROUND)
    _prepare_dark_axis(axis)
    axis.plot([0, 1], [0, 1], linestyle="--", color=COLOR_LIGHT_GRAY, linewidth=1.2, label="Ideal")

    palette = [COLOR_CYAN, COLOR_GOLD, COLOR_TEAL, COLOR_GREEN, COLOR_RED]
    if calibration_by_group.empty:
        axis.text(0.5, 0.5, "Calibration evidence unavailable", ha="center", va="center", color=COLOR_WHITE)
    else:
        for index, horizon in enumerate(sorted(calibration_by_group["Horizon"].unique())):
            subset = calibration_by_group[
                calibration_by_group["Horizon"] == horizon
            ].sort_values("Calibration_Group_Order")
            axis.plot(
                subset["Predicted_Survival"],
                subset["Observed_KM_Survival"],
                marker="o",
                linewidth=1.6,
                markersize=4,
                color=palette[index % len(palette)],
                label=f"{horizon:g} {config.time_unit_label}",
            )

    axis.set_xlim(0.0, 1.0)
    axis.set_ylim(0.0, 1.0)
    axis.set_aspect("equal", adjustable="box")
    axis.set_title("Out-of-Fold Survival Calibration by Risk Group")
    axis.set_xlabel("Mean Predicted Survival")
    axis.set_ylabel("Observed Kaplan-Meier Survival")
    legend = axis.legend(facecolor=COLOR_PANEL, edgecolor=COLOR_MID_GRAY, fontsize=8)
    for item in legend.get_texts():
        item.set_color(COLOR_WHITE)
    figure.tight_layout()
    figure.savefig(output_path, dpi=180, bbox_inches="tight", facecolor=figure.get_facecolor())
    plt.close(figure)


def plot_ph_assumption_diagnostics(
    model_artifacts: CoxModelArtifacts,
    output_path: Path,
) -> None:
    review_features = model_artifacts.ph_test_results.loc[
        model_artifacts.ph_test_results["Raw_PH_Status"].eq("REVIEW"), "Feature"
    ].astype(str).tolist()

    if not review_features or model_artifacts.ph_residuals.empty:
        figure, axis = plt.subplots(figsize=(8.5, 4.5), facecolor=COLOR_BACKGROUND)
        _prepare_dark_axis(axis)
        axis.axis("off")
        axis.text(
            0.5,
            0.55,
            "No feature-level raw PH diagnostic required a residual plot.",
            ha="center",
            va="center",
            color=COLOR_WHITE,
            fontsize=15,
        )
        axis.text(
            0.5,
            0.40,
            "Review proportional_hazards_test.csv for raw and adjusted p-values.",
            ha="center",
            va="center",
            color=COLOR_LIGHT_GRAY,
            fontsize=11,
        )
        figure.tight_layout()
        figure.savefig(output_path, dpi=180, bbox_inches="tight", facecolor=figure.get_facecolor())
        plt.close(figure)
        return

    plotted = [feature for feature in review_features if feature in set(model_artifacts.ph_residuals.get("Feature", []))]
    if not plotted:
        plotted = review_features[:1]
    plotted = plotted[:4]
    figure, axes = plt.subplots(
        len(plotted),
        1,
        figsize=(9.0, 3.2 * len(plotted)),
        facecolor=COLOR_BACKGROUND,
        squeeze=False,
    )

    for axis, feature in zip(axes.flatten(), plotted):
        _prepare_dark_axis(axis)
        subset = model_artifacts.ph_residuals[
            model_artifacts.ph_residuals["Feature"] == feature
        ].dropna(subset=["Event_Time", "Scaled_Schoenfeld_Residual"])
        axis.scatter(
            subset["Event_Time"],
            subset["Scaled_Schoenfeld_Residual"],
            s=10,
            alpha=0.25,
            color=COLOR_CYAN,
        )
        if len(subset) >= 3 and subset["Event_Time"].nunique() > 1:
            coefficients = np.polyfit(
                subset["Event_Time"].to_numpy(dtype=float),
                subset["Scaled_Schoenfeld_Residual"].to_numpy(dtype=float),
                deg=1,
            )
            x_values = np.linspace(subset["Event_Time"].min(), subset["Event_Time"].max(), 100)
            axis.plot(x_values, np.polyval(coefficients, x_values), color=COLOR_GOLD, linewidth=2)
        axis.axhline(0.0, color=COLOR_LIGHT_GRAY, linewidth=1)
        axis.set_title(f"Scaled Schoenfeld Residual Review — {feature}")
        axis.set_xlabel("Observed Event Time")
        axis.set_ylabel("Scaled Residual")

    figure.tight_layout()
    figure.savefig(output_path, dpi=180, bbox_inches="tight", facecolor=figure.get_facecolor())
    plt.close(figure)


def generate_evidence_plots(
    run_directory: Path,
    config: FrameworkConfig,
    personas: PersonaArtifacts,
    model: CoxModelArtifacts,
    risk: RiskArtifacts,
    scenarios: ScenarioArtifacts,
) -> PlotArtifacts:
    chart_directory = run_directory / "charts"
    persona_path = chart_directory / "persona_kaplan_meier.png"
    risk_path = chart_directory / "risk_tier_kaplan_meier.png"
    scenario_path = chart_directory / "baseline_vs_scenario_survival.png"
    comparison_path = chart_directory / "scenario_hazard_reduction.png"
    calibration_path = chart_directory / "calibration_at_horizons.png"
    ph_path = chart_directory / "ph_assumption_diagnostics.png"

    plot_persona_kaplan_meier(personas.scored_data, config, persona_path)
    plot_risk_tier_kaplan_meier(risk.scored_data, config, risk_path)
    plot_predicted_survival_comparison(
        scenarios.predicted_survival_curves,
        scenarios.summary,
        config,
        scenario_path,
    )
    plot_scenario_hazard_reduction(scenarios.summary, comparison_path)
    plot_calibration_evidence(model.calibration_by_risk_group, config, calibration_path)
    plot_ph_assumption_diagnostics(model, ph_path)

    return PlotArtifacts(
        persona_km=persona_path,
        risk_tier_km=risk_path,
        baseline_vs_scenario=scenario_path,
        scenario_comparison=comparison_path,
        calibration=calibration_path,
        ph_diagnostics=ph_path,
    )

# ==================================================================================================
# EXECUTIVE INTERPRETATION
# ==================================================================================================

def build_executive_narrative(
    config: FrameworkConfig,
    model: CoxModelArtifacts,
    risk: RiskArtifacts,
    scenarios: ScenarioArtifacts,
    acceptance_status: str,
) -> str:
    metadata = model.fit_metadata.iloc[0]
    event_rate = float(metadata["Event_Rate"])
    apparent_concordance = float(metadata["Apparent_Concordance_Index"])
    cv_mean = float(metadata["CV_Mean_Concordance"])
    cv_sd = float(metadata["CV_SD_Concordance"])
    target_count = int(len(risk.target_cohort))
    ph_status = str(metadata["PH_Global_Status"])

    improvement_rows = scenarios.summary[
        scenarios.summary["Scenario_Type"].eq("IMPROVEMENT")
    ]
    if improvement_rows.empty:
        scenario_text = "No favorable intervention scenarios were configured for this run."
    else:
        best = improvement_rows.sort_values("Improvement_Rank").iloc[0]
        scenario_text = (
            f"Among the illustrative improvement scenarios, '{best['Scenario_Name']}' produced "
            f"the largest modeled movement for the same {target_count:,}-record target cohort: "
            f"mean partial hazard declined by {best['Relative_Hazard_Reduction']:.2%}, and modeled "
            f"survival at {config.evaluation_horizon:g} {config.time_unit_label} increased by "
            f"{best['Survival_Uplift_Percentage_Points']:.2f} percentage points."
        )

    return (
        f"The {config.project_name} run analyzed {len(risk.scored_data):,} records with an observed "
        f"event rate of {event_rate:.2%}. The regularized CoxPH model produced an apparent "
        f"development-sample concordance index of {apparent_concordance:.3f}; deterministic "
        f"{config.cross_validation_folds}-fold validation produced a mean concordance of "
        f"{cv_mean:.3f} (SD {cv_sd:.3f}). The baseline {risk.target_label.lower()} rule selected "
        f"{target_count:,} records for same-cohort sensitivity testing. {scenario_text} The run "
        f"concluded {acceptance_status}; proportional-hazards evidence was classified {ph_status}. "
        f"These outputs are model-based sensitivity evidence from demonstration data, not causal "
        f"treatment effects, production retention forecasts, or realized financial outcomes."
    )


# ==================================================================================================
# EXECUTIVE POWERPOINT
# ==================================================================================================

def _ppt_set_background(slide: Any, rgb: RGBColor = RGBColor(7, 19, 26)) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _ppt_add_title(slide: Any, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(244, 245, 246)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.58), Inches(0.85), Inches(12.0), Inches(0.4))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.name = "Arial"
        subtitle_frame.paragraphs[0].font.size = Pt(12)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(32, 215, 229)


def _ppt_add_footer(slide: Any) -> None:
    footer_box = slide.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.25))
    frame = footer_box.text_frame
    frame.text = "Public demonstration data | Illustrative scenarios | Not a production retention forecast"
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = "Arial"
    paragraph.font.size = Pt(8)
    paragraph.font.color.rgb = RGBColor(113, 128, 138)


def _ppt_add_picture_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    subtitle: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_set_background(slide)
    _ppt_add_title(slide, title, subtitle)
    # Height is fixed to prevent charts from extending beneath the footer on a 16:9 slide.
    slide.shapes.add_picture(
        str(image_path),
        Inches(2.10),
        Inches(1.30),
        height=Inches(5.45),
    )
    _ppt_add_footer(slide)


def _ppt_add_kpi_card(
    slide: Any,
    left: float,
    top: float,
    width: float,
    value: str,
    label: str,
    accent: RGBColor,
) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1.05))
    frame = shape.text_frame
    frame.word_wrap = True
    frame.text = value
    value_paragraph = frame.paragraphs[0]
    value_paragraph.alignment = PP_ALIGN.CENTER
    value_paragraph.font.name = "Arial"
    value_paragraph.font.size = Pt(25)
    value_paragraph.font.bold = True
    value_paragraph.font.color.rgb = accent
    label_paragraph = frame.add_paragraph()
    label_paragraph.text = label
    label_paragraph.alignment = PP_ALIGN.CENTER
    label_paragraph.font.name = "Arial"
    label_paragraph.font.size = Pt(9)
    label_paragraph.font.color.rgb = RGBColor(199, 208, 213)


def _ppt_add_scenario_table(slide: Any, summary: pd.DataFrame) -> None:
    display_columns = [
        "Scenario_Name",
        "Scenario_Type",
        "Relative_Hazard_Reduction",
        "Survival_Uplift_Percentage_Points",
        "Modeled_Direction",
        "Scenario_Rank",
    ]
    table_df = summary.loc[:, display_columns].copy()
    table_df["Relative_Hazard_Reduction"] = table_df["Relative_Hazard_Reduction"].map(
        lambda value: f"{value:.1%}"
    )
    table_df["Survival_Uplift_Percentage_Points"] = table_df[
        "Survival_Uplift_Percentage_Points"
    ].map(lambda value: f"{value:.1f} pp")
    table_df = table_df.rename(
        columns={
            "Scenario_Name": "Scenario",
            "Scenario_Type": "Type",
            "Relative_Hazard_Reduction": "Hazard Reduction",
            "Survival_Uplift_Percentage_Points": "Survival Uplift",
            "Modeled_Direction": "Direction",
            "Scenario_Rank": "Rank",
        }
    )

    rows, columns = len(table_df) + 1, len(table_df.columns)
    shape = slide.shapes.add_table(
        rows,
        columns,
        Inches(0.45),
        Inches(1.45),
        Inches(12.45),
        Inches(4.95),
    )
    table = shape.table
    # Scenario name receives more horizontal space than compact metric columns.
    widths = [4.2, 1.25, 1.75, 1.55, 1.35, 0.65]
    for index, width in enumerate(widths):
        table.columns[index].width = Inches(width)

    for col_idx, column in enumerate(table_df.columns):
        table.cell(0, col_idx).text = column
    for row_idx, (_, row) in enumerate(table_df.iterrows(), start=1):
        for col_idx, value in enumerate(row.tolist()):
            table.cell(row_idx, col_idx).text = str(value)

    for row_idx in range(rows):
        for col_idx in range(columns):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                RGBColor(12, 32, 42) if row_idx else RGBColor(20, 73, 88)
            )
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.size = Pt(8.2)
                paragraph.font.color.rgb = RGBColor(244, 245, 246)
                if row_idx == 0:
                    paragraph.font.bold = True


def export_executive_powerpoint(
    output_path: Path,
    config: FrameworkConfig,
    model: CoxModelArtifacts,
    personas: PersonaArtifacts,
    risk: RiskArtifacts,
    scenarios: ScenarioArtifacts,
    plots: PlotArtifacts,
    narrative: str,
    acceptance_status: str,
) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 — Framework thesis.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_set_background(slide)
    _ppt_add_title(
        slide,
        "Survival Strategy Framework",
        "From time-to-event data to governed intervention evidence",
    )
    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(11.6), Inches(4.5))
    text_frame = body.text_frame
    text_frame.word_wrap = True
    bullets = [
        "Separate persona discovery from CoxPH risk modeling.",
        "Validate discrimination with stratified K-fold out-of-sample scoring.",
        "Select a model-derived top-quartile target cohort.",
        "Apply control, favorable, and adverse scenarios to the same governed IDs.",
        "Rebuild interactions and squared terms before re-scoring.",
        "Archive calibration, PH review, scenario controls, and stakeholder evidence.",
    ]
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(199, 208, 213)
        paragraph.space_after = Pt(10)
    _ppt_add_footer(slide)

    _ppt_add_picture_slide(
        prs,
        "Persona Lifecycle Evidence",
        plots.persona_km,
        "Observed Kaplan-Meier curves; personas are descriptive segments, not CoxPH model inputs.",
    )
    _ppt_add_picture_slide(
        prs,
        "Model-Derived Risk Tiers",
        plots.risk_tier_km,
        "Observed Kaplan-Meier evidence across baseline partial-hazard tiers.",
    )
    _ppt_add_picture_slide(
        prs,
        "Same-Cohort Predicted Survival",
        plots.baseline_vs_scenario,
        "CoxPH-predicted survival for identical target IDs under control, favorable, and stress scenarios.",
    )

    # Slide 5 — Scenario comparison table only. Splitting the chart prevents clipping.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_set_background(slide)
    _ppt_add_title(
        slide,
        "Scenario Comparison",
        f"Modeled sensitivity at {config.evaluation_horizon:g} {config.time_unit_label}; survival uplift shown in percentage points",
    )
    _ppt_add_scenario_table(slide, scenarios.summary)
    _ppt_add_footer(slide)

    # Slide 6 — Hazard movement chart only.
    _ppt_add_picture_slide(
        prs,
        "Relative Hazard Movement",
        plots.scenario_comparison,
        "Positive values indicate modeled hazard reduction; negative values indicate adverse stress movement.",
    )

    # Slide 7 — Cross-validation and calibration evidence.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_set_background(slide)
    _ppt_add_title(
        slide,
        "Model Validation Snapshot",
        "Out-of-fold discrimination, calibration, and proportional-hazards review",
    )
    slide.shapes.add_picture(
        str(plots.calibration),
        Inches(0.55),
        Inches(1.35),
        height=Inches(4.75),
    )
    metadata = model.fit_metadata.iloc[0]
    _ppt_add_kpi_card(
        slide,
        8.15,
        1.55,
        2.15,
        f"{metadata['Apparent_Concordance_Index']:.3f}",
        "Apparent concordance",
        RGBColor(217, 164, 65),
    )
    _ppt_add_kpi_card(
        slide,
        10.35,
        1.55,
        2.15,
        f"{metadata['CV_Mean_Concordance']:.3f}",
        f"{config.cross_validation_folds}-fold CV mean",
        RGBColor(32, 215, 229),
    )
    _ppt_add_kpi_card(
        slide,
        8.15,
        3.15,
        2.15,
        f"{metadata['Integrated_Brier_Score']:.3f}",
        "Integrated Brier score",
        RGBColor(86, 194, 113),
    )
    _ppt_add_kpi_card(
        slide,
        10.35,
        3.15,
        2.15,
        str(metadata["PH_Global_Status"]),
        "PH diagnostic status",
        RGBColor(217, 164, 65),
    )
    _ppt_add_footer(slide)

    # Slide 8 — Executive interpretation and explicit limits.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_set_background(slide)
    _ppt_add_title(slide, "Executive Interpretation", "Evidence, review status, and boundaries")

    metadata = model.fit_metadata.iloc[0]
    improvement_rows = scenarios.summary[scenarios.summary["Scenario_Type"].eq("IMPROVEMENT")]
    eligible_rows = (
        improvement_rows
        if not improvement_rows.empty
        else scenarios.summary[~scenarios.summary["Scenario_Type"].eq("CONTROL")]
    )
    best = eligible_rows.sort_values(
        "Relative_Hazard_Reduction", ascending=False
    ).iloc[0]
    _ppt_add_kpi_card(slide, 0.65, 1.30, 2.35, f"{len(risk.scored_data):,}", "Synthetic records", RGBColor(32, 215, 229))
    _ppt_add_kpi_card(slide, 3.05, 1.30, 2.35, f"{len(risk.target_cohort):,}", "Target-cohort records", RGBColor(32, 215, 229))
    _ppt_add_kpi_card(slide, 5.45, 1.30, 2.35, f"{best['Relative_Hazard_Reduction']:.1%}", "Best modeled hazard reduction", RGBColor(86, 194, 113))
    _ppt_add_kpi_card(slide, 7.85, 1.30, 2.35, f"{best['Survival_Uplift_Percentage_Points']:.1f} pp", f"Survival uplift at {config.evaluation_horizon:g} {config.time_unit_label}", RGBColor(86, 194, 113))
    _ppt_add_kpi_card(slide, 10.25, 1.30, 2.35, acceptance_status, "Run acceptance", RGBColor(217, 164, 65))

    narrative_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.75), Inches(11.8), Inches(1.65))
    narrative_frame = narrative_box.text_frame
    narrative_frame.word_wrap = True
    narrative_frame.text = narrative
    narrative_frame.paragraphs[0].font.name = "Arial"
    narrative_frame.paragraphs[0].font.size = Pt(13)
    narrative_frame.paragraphs[0].font.color.rgb = RGBColor(244, 245, 246)

    boundary_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.70), Inches(11.8), Inches(1.35))
    boundary_frame = boundary_box.text_frame
    boundaries = [
        "Modeled sensitivity — not causal treatment effect.",
        "Synthetic demonstration data — no PII.",
        "No booked revenue, realized ROI, or production forecast is claimed.",
        "Production use requires domain calibration, monitoring, and model-risk approval.",
    ]
    for index, boundary in enumerate(boundaries):
        paragraph = boundary_frame.paragraphs[0] if index == 0 else boundary_frame.add_paragraph()
        paragraph.text = boundary
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(217, 164, 65)
        paragraph.space_after = Pt(5)
    _ppt_add_footer(slide)

    prs.save(output_path)

# ==================================================================================================
# TECHNICAL PDF
# ==================================================================================================

def dataframe_for_display(df: pd.DataFrame, max_rows: int = 30) -> pd.DataFrame:
    """Formats evidence tables for human review without altering archived raw CSV values."""
    display = df.head(max_rows).copy()
    for column in display.columns:
        lower = column.lower()
        if pd.api.types.is_float_dtype(display[column]):
            if "percentage_points" in lower or lower.endswith("_pp"):
                display[column] = display[column].map(
                    lambda value: "" if pd.isna(value) else f"{value:.2f} pp"
                )
            elif "p_value" in lower or lower in {"p", "alpha"}:
                display[column] = display[column].map(
                    lambda value: "" if pd.isna(value) else f"{value:.4g}"
                )
            elif any(
                token in lower
                for token in [
                    "rate",
                    "reduction",
                    "uplift",
                    "survival",
                    "relative_change",
                    "overlap_rate",
                    "calibration_difference",
                    "calibration_error",
                ]
            ):
                display[column] = display[column].map(
                    lambda value: "" if pd.isna(value) else f"{value:.2%}"
                )
            else:
                display[column] = display[column].map(
                    lambda value: "" if pd.isna(value) else f"{value:,.4f}"
                )
    return display


def reportlab_table_from_df(
    df: pd.DataFrame,
    max_rows: int = 30,
    font_size: float = 6.5,
) -> Table:
    display = dataframe_for_display(df, max_rows=max_rows)
    data: list[list[Any]] = [
        [
            Paragraph(
                str(column),
                ParagraphStyle("Header", fontSize=font_size, textColor=colors.white),
            )
            for column in display.columns
        ]
    ]
    for _, row in display.iterrows():
        data.append(
            [
                Paragraph(
                    "" if pd.isna(value) else str(value),
                    ParagraphStyle("Cell", fontSize=font_size, leading=font_size + 1.2),
                )
                for value in row.tolist()
            ]
        )

    usable_width = landscape(letter)[0] - 0.7 * inch
    column_width = usable_width / max(1, len(display.columns))
    table = Table(data, repeatRows=1, colWidths=[column_width] * len(display.columns))
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#14566A")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#F4F6F7")),
                ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#87949C")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 3),
                ("RIGHTPADDING", (0, 0), (-1, -1), 3),
                ("TOPPADDING", (0, 0), (-1, -1), 3),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ]
        )
    )
    return table


def export_technical_pdf(
    output_path: Path,
    config: FrameworkConfig,
    validation: pd.DataFrame,
    acceptance_checks: pd.DataFrame,
    reproducibility_checks: pd.DataFrame,
    acceptance_status: str,
    model: CoxModelArtifacts,
    personas: PersonaArtifacts,
    risk: RiskArtifacts,
    scenarios: ScenarioArtifacts,
    plots: PlotArtifacts,
    narrative: str,
) -> None:
    document = SimpleDocTemplate(
        str(output_path),
        pagesize=landscape(letter),
        leftMargin=0.35 * inch,
        rightMargin=0.35 * inch,
        topMargin=0.35 * inch,
        bottomMargin=0.35 * inch,
    )
    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="CenteredTitle",
            parent=styles["Title"],
            alignment=TA_CENTER,
            textColor=colors.HexColor("#14566A"),
        )
    )
    styles.add(
        ParagraphStyle(
            name="EvidenceHeading",
            parent=styles["Heading2"],
            textColor=colors.HexColor("#14566A"),
            spaceBefore=8,
            spaceAfter=6,
        )
    )
    styles.add(
        ParagraphStyle(
            name="StatusLine",
            parent=styles["Heading2"],
            alignment=TA_CENTER,
            textColor=(
                colors.HexColor("#8A6D1D")
                if acceptance_status == "PASS_WITH_REVIEW"
                else (
                    colors.HexColor("#A32929")
                    if acceptance_status == "FAIL"
                    else colors.HexColor("#2B7A3D")
                )
            ),
        )
    )
    story: list[Any] = []

    story.append(Paragraph("Technical Model Evidence", styles["CenteredTitle"]))
    story.append(Paragraph(config.project_name, styles["Heading2"]))
    story.append(Paragraph(f"Run conclusion: {acceptance_status}", styles["StatusLine"]))
    story.append(Paragraph(narrative, styles["BodyText"]))
    story.append(Spacer(1, 10))

    story.append(Paragraph("1. Acceptance Checks", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(acceptance_checks, max_rows=50, font_size=7))

    story.append(Paragraph("2. Input Validation", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(validation, max_rows=50, font_size=7))

    story.append(PageBreak())
    story.append(Paragraph("3. Model Fit Metadata", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(model.fit_metadata, max_rows=10, font_size=6.5))

    story.append(Paragraph("4. Stratified K-Fold Concordance", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(model.cross_validation_results, max_rows=20, font_size=7))

    story.append(Paragraph("5. Out-of-Fold Calibration Metrics", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(model.calibration_metrics, max_rows=20, font_size=6.5))
    story.append(Spacer(1, 7))
    story.append(RLImage(str(plots.calibration), width=4.7 * inch, height=3.8 * inch))

    story.append(PageBreak())
    story.append(Paragraph("6. CoxPH Model Summary", styles["EvidenceHeading"]))
    model_columns = [
        column
        for column in [
            "Feature",
            "Coefficient",
            "Hazard_Ratio",
            "Hazard_Ratio_CI_Lower",
            "Hazard_Ratio_CI_Upper",
            "P_Value",
            "Interpretation_Unit",
        ]
        if column in model.model_summary.columns
    ]
    story.append(reportlab_table_from_df(model.model_summary[model_columns], max_rows=50, font_size=6.3))

    story.append(PageBreak())
    story.append(Paragraph("7. Proportional-Hazards Tests", styles["EvidenceHeading"]))
    ph_columns = [
        column
        for column in [
            "Feature",
            "Test_Statistic",
            "P_Value",
            "P_Value_Holm",
            "P_Value_BH",
            "Raw_PH_Status",
            "Holm_PH_Status",
            "PH_Assumption_Status",
            "Review_Rationale",
        ]
        if column in model.ph_test_results.columns
    ]
    story.append(reportlab_table_from_df(model.ph_test_results[ph_columns], max_rows=50, font_size=6.0))
    story.append(Paragraph("8. PH Sensitivity Review", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(model.ph_sensitivity_results, max_rows=30, font_size=6.2))
    story.append(Spacer(1, 7))
    story.append(RLImage(str(plots.ph_diagnostics), width=5.4 * inch, height=3.0 * inch))

    story.append(PageBreak())
    story.append(Paragraph("9. Persona Quality and Stability", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(personas.quality_metrics, max_rows=10, font_size=6.5))
    story.append(reportlab_table_from_df(personas.stability_results, max_rows=20, font_size=6.5))
    story.append(Paragraph("10. Persona Profiles", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(personas.profiles, max_rows=20, font_size=6.2))
    story.append(Spacer(1, 7))
    story.append(RLImage(str(plots.persona_km), width=5.1 * inch, height=3.05 * inch))

    story.append(PageBreak())
    story.append(Paragraph("11. Risk-Tier Profiles", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(risk.risk_tier_profiles, max_rows=20, font_size=6.2))
    story.append(Spacer(1, 8))
    story.append(RLImage(str(plots.risk_tier_km), width=5.2 * inch, height=3.1 * inch))

    story.append(PageBreak())
    story.append(Paragraph("12. Scenario Results", styles["EvidenceHeading"]))
    scenario_columns = [
        "Scenario_Rank",
        "Improvement_Rank",
        "Scenario_Name",
        "Scenario_Type",
        "Target_Cohort_Count",
        "Relative_Hazard_Reduction",
        "Baseline_Survival_At_Horizon",
        "Scenario_Survival_At_Horizon",
        "Survival_Uplift_Percentage_Points",
        "Modeled_Direction",
        "Direction_Check_Status",
    ]
    story.append(reportlab_table_from_df(scenarios.summary[scenario_columns], max_rows=30, font_size=5.9))
    story.append(Spacer(1, 8))
    story.append(RLImage(str(plots.baseline_vs_scenario), width=5.1 * inch, height=3.05 * inch))

    story.append(PageBreak())
    story.append(Paragraph("13. Scenario Change Audit", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(scenarios.scenario_change_audit, max_rows=80, font_size=5.6))
    story.append(Spacer(1, 8))
    story.append(RLImage(str(plots.scenario_comparison), width=5.2 * inch, height=3.1 * inch))

    story.append(PageBreak())
    story.append(Paragraph("14. Dependency Synchronization Audit", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(scenarios.dependency_audit, max_rows=100, font_size=5.8))

    story.append(Paragraph("15. Reproducibility Checks", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(reproducibility_checks, max_rows=30, font_size=7))

    story.append(Paragraph("16. Interpretation Boundaries", styles["EvidenceHeading"]))
    boundaries = [
        "The analysis uses synthetic or public demonstration data.",
        "Scenario movement is modeled sensitivity, not causal treatment effect.",
        "Partial hazard is relative risk, not an event probability.",
        "Apparent concordance is in-sample; cross-validation is reported separately.",
        "Out-of-fold calibration does not replace external or temporal validation.",
        "Predicted survival requires domain validation before production use.",
        "No realized revenue, ROI, or defended LTV is claimed.",
    ]
    for boundary in boundaries:
        story.append(Paragraph(f"• {boundary}", styles["BodyText"]))

    document.build(story)

# ==================================================================================================
# RUN ARCHIVE, ACCEPTANCE CHECKS, REPRODUCIBILITY, AND MANIFEST
# ==================================================================================================

def verify_substantive_reproducibility(
    df: pd.DataFrame,
    config: FrameworkConfig,
    scenarios_config: Sequence[ScenarioDefinition],
    personas: PersonaArtifacts,
    model: CoxModelArtifacts,
    risk: RiskArtifacts,
    scenarios: ScenarioArtifacts,
) -> pd.DataFrame:
    """Re-runs the analytical core and compares substantive outputs.

    This is intentionally stronger than confirming that the same input fingerprint was reused. It
    verifies persona assignments, coefficients, cross-validation, calibration, target IDs, scenario
    metrics, and predicted curves under the same deterministic configuration.
    """
    if not config.verify_reproducibility:
        return pd.DataFrame(
            [{
                "Check": "Substantive reproducibility verification",
                "Status": "NOT_RUN",
                "Evidence": "Disabled by framework configuration.",
            }]
        )

    second_personas = fit_personas(df, config)
    second_model = fit_cox_model(df, config)
    second_risk = score_and_stratify_population(df, second_personas, second_model, config)
    second_scenarios = simulate_scenarios(
        df,
        config,
        scenarios_config,
        second_model,
        second_risk,
    )
    second_model.ph_sensitivity_results = run_ph_sensitivity_analysis(
        df,
        config,
        second_model,
        second_risk,
        scenarios_config=scenarios_config,
        primary_scenario_summary=second_scenarios.summary,
    )

    rows: list[dict[str, str]] = []

    def record(name: str, passed: bool, evidence: str) -> None:
        rows.append(
            {
                "Check": name,
                "Status": "PASS" if passed else "FAIL",
                "Evidence": evidence,
            }
        )

    record(
        "Persona assignments reproducible",
        personas.scored_data["Segment_ID"].equals(second_personas.scored_data["Segment_ID"]),
        "Segment_ID assignments were compared row-for-row under the same seed and configuration.",
    )

    first_coefficients = model.model_summary.set_index("Feature")["Coefficient"].sort_index()
    second_coefficients = second_model.model_summary.set_index("Feature")["Coefficient"].sort_index()
    coefficient_difference = float(
        np.max(np.abs(first_coefficients.to_numpy() - second_coefficients.to_numpy()))
    )
    record(
        "CoxPH coefficients reproducible",
        coefficient_difference <= 1e-10,
        f"Maximum absolute coefficient difference: {coefficient_difference:.3e}.",
    )

    first_cv = model.cross_validation_results.sort_values("Fold")["Concordance_Index"].to_numpy(dtype=float)
    second_cv = second_model.cross_validation_results.sort_values("Fold")["Concordance_Index"].to_numpy(dtype=float)
    cv_difference = float(np.nanmax(np.abs(first_cv - second_cv)))
    record(
        "Cross-validation evidence reproducible",
        cv_difference <= 1e-10,
        f"Maximum absolute fold concordance difference: {cv_difference:.3e}.",
    )

    calibration_columns = [
        "Horizon",
        "Overall_Predicted_Survival",
        "Overall_Observed_KM_Survival",
        "IPCW_Brier_Score",
    ]
    first_calibration = model.calibration_metrics[calibration_columns].sort_values("Horizon")
    second_calibration = second_model.calibration_metrics[calibration_columns].sort_values("Horizon")
    calibration_difference = float(
        np.nanmax(
            np.abs(
                first_calibration.select_dtypes(include=[np.number]).to_numpy()
                - second_calibration.select_dtypes(include=[np.number]).to_numpy()
            )
        )
    )
    record(
        "Calibration evidence reproducible",
        calibration_difference <= 1e-10,
        f"Maximum absolute calibration difference: {calibration_difference:.3e}.",
    )

    ph_columns = [
        column
        for column in ["P_Value", "P_Value_Holm", "P_Value_BH"]
        if column in model.ph_test_results.columns
        and column in second_model.ph_test_results.columns
    ]
    first_ph = model.ph_test_results.sort_values("Feature").reset_index(drop=True)
    second_ph = second_model.ph_test_results.sort_values("Feature").reset_index(drop=True)
    if ph_columns and len(first_ph) == len(second_ph):
        ph_difference = float(
            np.nanmax(
                np.abs(
                    first_ph[ph_columns].to_numpy(dtype=float)
                    - second_ph[ph_columns].to_numpy(dtype=float)
                )
            )
        )
        ph_status_match = first_ph["PH_Assumption_Status"].equals(
            second_ph["PH_Assumption_Status"]
        )
    else:
        ph_difference = np.inf
        ph_status_match = False
    record(
        "PH diagnostic evidence reproducible",
        ph_difference <= 1e-10 and ph_status_match,
        (
            f"Maximum absolute raw/adjusted PH p-value difference: {ph_difference:.3e}; "
            f"status match={ph_status_match}."
        ),
    )

    first_ids = set(risk.target_cohort[config.id_col].astype(str))
    second_ids = set(second_risk.target_cohort[config.id_col].astype(str))
    record(
        "Target cohort reproducible",
        first_ids == second_ids,
        f"Compared {len(first_ids):,} baseline target IDs with the repeated run.",
    )

    metric_columns = [
        "Baseline_Mean_Partial_Hazard",
        "Scenario_Mean_Partial_Hazard",
        "Relative_Hazard_Reduction",
        "Baseline_Survival_At_Horizon",
        "Scenario_Survival_At_Horizon",
        "Survival_Probability_Uplift",
    ]
    first_summary = scenarios.summary.sort_values("Scenario_Name").set_index("Scenario_Name")
    second_summary = second_scenarios.summary.sort_values("Scenario_Name").set_index("Scenario_Name")
    scenario_difference = float(
        np.nanmax(
            np.abs(
                first_summary[metric_columns].to_numpy(dtype=float)
                - second_summary[metric_columns].to_numpy(dtype=float)
            )
        )
    )
    directions_match = first_summary["Modeled_Direction"].equals(
        second_summary["Modeled_Direction"]
    )
    record(
        "Scenario evidence reproducible",
        scenario_difference <= 1e-10 and directions_match,
        f"Maximum absolute scenario metric difference: {scenario_difference:.3e}; directions match={directions_match}.",
    )

    first_curves = scenarios.predicted_survival_curves.sort_values("Time")
    second_curves = second_scenarios.predicted_survival_curves.sort_values("Time")
    curve_difference = float(
        np.nanmax(
            np.abs(
                first_curves.select_dtypes(include=[np.number]).to_numpy()
                - second_curves.select_dtypes(include=[np.number]).to_numpy()
            )
        )
    )
    record(
        "Predicted survival curves reproducible",
        curve_difference <= 1e-10,
        f"Maximum absolute curve difference: {curve_difference:.3e}.",
    )

    first_sensitivity = model.ph_sensitivity_results.sort_values("Feature").reset_index(drop=True)
    second_sensitivity = second_model.ph_sensitivity_results.sort_values("Feature").reset_index(drop=True)
    sensitivity_keys = [
        column
        for column in [
            "Feature",
            "Fit_Status",
            "Scenario_Rank_Status",
            "Primary_Best_Improvement",
            "Sensitivity_Best_Improvement",
        ]
        if column in first_sensitivity.columns and column in second_sensitivity.columns
    ]
    sensitivity_numeric = [
        column
        for column in [
            "Apparent_Concordance_Index",
            "Concordance_Delta",
            "Target_Cohort_Overlap_Rate",
            "Target_Cohort_Jaccard",
            "Scenario_Rank_Spearman",
        ]
        if column in first_sensitivity.columns and column in second_sensitivity.columns
    ]
    sensitivity_string_match = (
        len(first_sensitivity) == len(second_sensitivity)
        and (
            first_sensitivity[sensitivity_keys].fillna("").astype(str).equals(
                second_sensitivity[sensitivity_keys].fillna("").astype(str)
            )
            if sensitivity_keys
            else True
        )
    )
    if sensitivity_numeric and len(first_sensitivity) == len(second_sensitivity):
        first_numeric = first_sensitivity[sensitivity_numeric].to_numpy(dtype=float)
        second_numeric = second_sensitivity[sensitivity_numeric].to_numpy(dtype=float)
        sensitivity_difference = float(
            np.nanmax(np.abs(first_numeric - second_numeric))
        )
        if np.isnan(sensitivity_difference):
            sensitivity_difference = 0.0
    else:
        sensitivity_difference = 0.0 if not sensitivity_numeric else np.inf
    record(
        "PH sensitivity evidence reproducible",
        sensitivity_string_match and sensitivity_difference <= 1e-10,
        (
            f"Maximum absolute PH sensitivity metric difference: "
            f"{sensitivity_difference:.3e}; categorical evidence match="
            f"{sensitivity_string_match}."
        ),
    )

    return pd.DataFrame(rows)


def run_acceptance_checks(
    df: pd.DataFrame,
    config: FrameworkConfig,
    personas: PersonaArtifacts,
    model: CoxModelArtifacts,
    risk: RiskArtifacts,
    scenarios: ScenarioArtifacts,
    reproducibility_checks: pd.DataFrame,
    report_paths: Sequence[Path] = (),
) -> pd.DataFrame:
    """Creates PASS / REVIEW / FAIL evidence for the complete framework run."""
    checks: list[dict[str, Any]] = []

    def add(name: str, status: str, evidence: str) -> None:
        checks.append({"Check": name, "Status": status, "Evidence": evidence})

    add(
        "Unique governed IDs",
        "PASS" if not df[config.id_col].duplicated().any() else "FAIL",
        f"{df[config.id_col].nunique():,} unique IDs across {len(df):,} rows.",
    )
    add(
        "Target cohort is non-empty",
        "PASS" if len(risk.target_cohort) > 0 else "FAIL",
        f"{len(risk.target_cohort):,} records selected at quantile {config.risk_quantile:.2f}.",
    )

    expected_ids = set(risk.target_cohort[config.id_col].astype(str))
    consistent_ids = True
    if not scenarios.target_scores.empty:
        for _, group in scenarios.target_scores.groupby("Scenario_Name"):
            if set(group[config.id_col].astype(str)) != expected_ids:
                consistent_ids = False
                break
    add(
        "Same target cohort across scenarios",
        "PASS" if consistent_ids else "FAIL",
        "Every scenario score table was compared with the identical baseline target IDs.",
    )

    dependency_pass = (
        scenarios.dependency_audit.empty
        or scenarios.dependency_audit["Validation_Status"].isin(["PASS", "NOT_APPLICABLE"]).all()
    )
    add(
        "Engineered dependencies synchronized",
        "PASS" if dependency_pass else "FAIL",
        "All interaction and squared-term post-rebuild checks were reviewed.",
    )

    change_audit_complete = (
        not scenarios.scenario_change_audit.empty
        and {
            "Records_Changed_Rate",
            "Lower_Bound_Hit_Rate",
            "Upper_Bound_Hit_Rate",
            "Mean_Signed_Change",
        }.issubset(scenarios.scenario_change_audit.columns)
    )
    add(
        "Scenario change audit completed",
        "PASS" if change_audit_complete else "FAIL",
        "Configured changes, realized movement, and bound-hit rates were archived.",
    )

    finite_metrics = True
    if not scenarios.summary.empty:
        metric_columns = [
            "Baseline_Mean_Partial_Hazard",
            "Scenario_Mean_Partial_Hazard",
            "Relative_Hazard_Reduction",
            "Baseline_Survival_At_Horizon",
            "Scenario_Survival_At_Horizon",
            "Survival_Probability_Uplift",
        ]
        finite_metrics = np.isfinite(
            scenarios.summary[metric_columns].to_numpy(dtype=float)
        ).all()
    add(
        "Scenario metrics are calculated and finite",
        "PASS" if finite_metrics else "FAIL",
        "No hard-coded uncertainty, ROI, revenue, or defended-value metric is used.",
    )

    curve_values = scenarios.predicted_survival_curves.drop(columns=["Time"])
    curve_bounds = curve_values.apply(lambda column: column.between(0.0, 1.0).all()).all()
    add(
        "Predicted survival remains within [0, 1]",
        "PASS" if curve_bounds else "FAIL",
        "All baseline and scenario survival probabilities satisfy probability bounds.",
    )

    first_row = scenarios.predicted_survival_curves.iloc[0]
    time_zero_pass = (
        abs(float(first_row["Time"])) <= config.neutral_scenario_tolerance
        and np.allclose(
            first_row.drop(labels=["Time"]).to_numpy(dtype=float),
            1.0,
            atol=config.neutral_scenario_tolerance,
            rtol=0.0,
        )
    )
    add(
        "Predicted survival begins at S(0)=1",
        "PASS" if time_zero_pass else "FAIL",
        "The first baseline and scenario survival values were tested at time zero.",
    )

    direction_pass = scenarios.summary["Direction_Check_Status"].eq("PASS").all()
    add(
        "Scenario direction matches configured expectation",
        "PASS" if direction_pass else "FAIL",
        "Control, favorable, stress, and technical scenarios were compared with expected direction.",
    )

    control = scenarios.summary[scenarios.summary["Scenario_Type"].eq("CONTROL")]
    if control.empty:
        add("No-change control scenario", "REVIEW", "No CONTROL scenario was configured.")
    else:
        neutral = (
            control["Relative_Hazard_Reduction"].abs().le(config.neutral_scenario_tolerance).all()
            and control["Survival_Probability_Uplift"].abs().le(config.neutral_scenario_tolerance).all()
        )
        add(
            "No-change control scenario",
            "PASS" if neutral else "FAIL",
            "Expected zero modeled movement within numerical tolerance.",
        )

    stress = scenarios.summary[scenarios.summary["Scenario_Type"].eq("STRESS")]
    if stress.empty:
        add("Adverse stress scenario", "REVIEW", "No STRESS scenario was configured.")
    else:
        adverse = (
            stress["Relative_Hazard_Reduction"].lt(0).all()
            and stress["Survival_Probability_Uplift"].lt(0).all()
        )
        add(
            "Adverse stress scenario",
            "PASS" if adverse else "FAIL",
            "Expected negative hazard-reduction and survival-uplift evidence.",
        )

    cv_pass = (
        model.cross_validation_results["Status"].eq("PASS").all()
        and np.isfinite(model.cross_validation_results["Concordance_Index"]).all()
    )
    add(
        "Cross-validation completed",
        "PASS" if cv_pass else "FAIL",
        f"{config.cross_validation_folds} stratified folds were required.",
    )

    calibration_status = str(model.fit_metadata.iloc[0]["Calibration_Status"])
    add(
        "Out-of-fold calibration evidence",
        "PASS" if calibration_status == "PASS" else "REVIEW",
        (
            f"Maximum group absolute calibration error: "
            f"{model.fit_metadata.iloc[0]['Max_Group_Absolute_Calibration_Error']:.2%}."
        ),
    )

    persona_status = str(personas.quality_metrics.iloc[0]["Persona_Quality_Status"])
    add(
        "Persona quality evidence",
        "PASS" if persona_status == "PASS" else "REVIEW",
        (
            f"Silhouette={personas.quality_metrics.iloc[0]['Silhouette_Score']:.3f}; "
            f"mean stability ARI={personas.quality_metrics.iloc[0]['Stability_Mean_ARI']:.3f}."
        ),
    )

    ph_status = str(model.fit_metadata.iloc[0]["PH_Global_Status"])
    add(
        "Proportional-hazards diagnostics",
        "PASS" if ph_status == "PASS" else "REVIEW",
        f"Global documented PH status: {ph_status}.",
    )

    sensitivity_statuses = model.ph_sensitivity_results.get(
        "Fit_Status", pd.Series(dtype=str)
    ).astype(str)
    ranking_statuses = model.ph_sensitivity_results.get(
        "Scenario_Rank_Status", pd.Series(dtype=str)
    ).astype(str)
    sensitivity_review = (
        sensitivity_statuses.isin(["FAIL", "NOT_RUN"]).any()
        or ranking_statuses.isin(["REVIEW", "NOT_RUN", "NOT_AVAILABLE"]).any()
    )
    add(
        "PH sensitivity review",
        "REVIEW" if sensitivity_review else "PASS",
        (
            "Low-cardinality raw PH review features were stratified where technically appropriate; "
            "target-cohort overlap and favorable-scenario ranking were compared."
        ),
    )

    if config.verify_reproducibility:
        reproducible = reproducibility_checks["Status"].eq("PASS").all()
        add(
            "Substantive reproducibility",
            "PASS" if reproducible else "FAIL",
            "Persona, model, validation, target-cohort, scenario, and curve outputs were repeated.",
        )
    else:
        add("Substantive reproducibility", "REVIEW", "Verification was disabled by configuration.")

    if report_paths:
        reports_exist = all(path.exists() and path.stat().st_size > 0 for path in report_paths)
        add(
            "Executive and technical reports created",
            "PASS" if reports_exist else "FAIL",
            " | ".join(path.name for path in report_paths),
        )

    return pd.DataFrame(checks)


def archive_run_tables(
    run_directory: Path,
    config: FrameworkConfig,
    scenarios_config: Sequence[ScenarioDefinition],
    input_df: pd.DataFrame,
    validation: pd.DataFrame,
    acceptance_checks: pd.DataFrame,
    reproducibility_checks: pd.DataFrame,
    personas: PersonaArtifacts,
    model: CoxModelArtifacts,
    risk: RiskArtifacts,
    scenarios: ScenarioArtifacts,
    narrative: str,
) -> None:
    write_json(run_directory / "framework_config.json", config)
    write_json(run_directory / "scenario_definitions.json", list(scenarios_config))

    validation.to_csv(run_directory / "input_validation.csv", index=False)
    acceptance_checks.to_csv(run_directory / "acceptance_checks.csv", index=False)
    reproducibility_checks.to_csv(run_directory / "reproducibility_checks.csv", index=False)
    if config.save_input_snapshot:
        input_df.to_csv(run_directory / "input_snapshot.csv", index=False)

    model.model_summary.to_csv(run_directory / "model_summary.csv", index=False)
    model.fit_metadata.to_csv(run_directory / "model_fit_metadata.csv", index=False)
    model.fit_warnings.to_csv(run_directory / "model_fit_warnings.csv", index=False)
    model.cross_validation_results.to_csv(run_directory / "cross_validation_results.csv", index=False)
    model.cross_validated_predictions.to_csv(
        run_directory / "cross_validated_predictions.csv", index=False
    )
    model.calibration_metrics.to_csv(run_directory / "calibration_metrics.csv", index=False)
    model.calibration_by_risk_group.to_csv(
        run_directory / "calibration_by_risk_group.csv", index=False
    )
    model.ph_test_results.to_csv(run_directory / "proportional_hazards_test.csv", index=False)
    model.ph_residuals.to_csv(run_directory / "scaled_schoenfeld_residuals.csv", index=False)
    model.ph_sensitivity_results.to_csv(
        run_directory / "ph_sensitivity_results.csv", index=False
    )

    persona_scaler_parameters = pd.DataFrame(
        {
            "Feature": list(config.segmentation_features),
            "Mean": personas.scaler.mean_,
            "Scale": personas.scaler.scale_,
        }
    )
    persona_scaler_parameters.to_csv(
        run_directory / "persona_scaler_parameters.csv", index=False
    )

    standardized_centroids = pd.DataFrame(
        personas.kmeans.cluster_centers_,
        columns=list(config.segmentation_features),
    )
    standardized_centroids.insert(
        0,
        "Persona",
        [personas.persona_name_map[index] for index in range(len(standardized_centroids))],
    )
    standardized_centroids.insert(0, "Segment_ID", np.arange(len(standardized_centroids)))
    standardized_centroids.to_csv(
        run_directory / "persona_centroids_standardized.csv", index=False
    )

    raw_centroids = pd.DataFrame(
        personas.scaler.inverse_transform(personas.kmeans.cluster_centers_),
        columns=list(config.segmentation_features),
    )
    raw_centroids.insert(
        0,
        "Persona",
        [personas.persona_name_map[index] for index in range(len(raw_centroids))],
    )
    raw_centroids.insert(0, "Segment_ID", np.arange(len(raw_centroids)))
    raw_centroids.to_csv(
        run_directory / "persona_centroids_raw_scale.csv", index=False
    )

    personas.quality_metrics.to_csv(
        run_directory / "persona_quality_metrics.csv", index=False
    )
    personas.stability_results.to_csv(
        run_directory / "persona_stability_results.csv", index=False
    )
    personas.profiles.to_csv(run_directory / "persona_profiles.csv", index=False)

    if model.model_scaler is not None:
        model_scaler_parameters = pd.DataFrame(
            {
                "Feature": list(config.model_features),
                "Mean": model.model_scaler.mean_,
                "Scale": model.model_scaler.scale_,
            }
        )
    else:
        model_scaler_parameters = pd.DataFrame(
            {
                "Feature": list(config.model_features),
                "Mean": np.nan,
                "Scale": np.nan,
                "Note": "Model features were not standardized.",
            }
        )
    model_scaler_parameters.to_csv(
        run_directory / "model_scaler_parameters.csv", index=False
    )

    baseline_survival = model.model.baseline_survival_.copy()
    if 0.0 not in baseline_survival.index:
        origin = pd.DataFrame(
            [[1.0] * baseline_survival.shape[1]],
            index=[0.0],
            columns=baseline_survival.columns,
        )
        baseline_survival = pd.concat([origin, baseline_survival]).sort_index()
    else:
        baseline_survival.loc[0.0, :] = 1.0
    baseline_survival.reset_index().to_csv(
        run_directory / "cox_baseline_survival.csv", index=False
    )

    baseline_cumulative_hazard = model.model.baseline_cumulative_hazard_.copy()
    if 0.0 not in baseline_cumulative_hazard.index:
        origin = pd.DataFrame(
            [[0.0] * baseline_cumulative_hazard.shape[1]],
            index=[0.0],
            columns=baseline_cumulative_hazard.columns,
        )
        baseline_cumulative_hazard = pd.concat(
            [origin, baseline_cumulative_hazard]
        ).sort_index()
    else:
        baseline_cumulative_hazard.loc[0.0, :] = 0.0
    baseline_cumulative_hazard.reset_index().to_csv(
        run_directory / "cox_baseline_cumulative_hazard.csv", index=False
    )

    risk.risk_tier_profiles.to_csv(run_directory / "risk_tier_profiles.csv", index=False)
    risk.target_cohort.to_csv(run_directory / "target_cohort.csv", index=False)

    scenarios.summary.to_csv(run_directory / "scenario_results.csv", index=False)
    scenarios.scenario_change_audit.to_csv(
        run_directory / "scenario_change_audit.csv", index=False
    )
    scenarios.target_scores.to_csv(run_directory / "scenario_target_scores.csv", index=False)
    scenarios.dependency_audit.to_csv(run_directory / "dependency_audit.csv", index=False)
    scenarios.predicted_survival_curves.to_csv(
        run_directory / "predicted_survival_curves.csv", index=False
    )

    (run_directory / "executive_narrative.txt").write_text(narrative, encoding="utf-8")


def build_manifest(run_directory: Path) -> dict[str, str]:
    return {
        path.relative_to(run_directory).as_posix(): str(path.resolve())
        for path in sorted(run_directory.rglob("*"))
        if path.is_file()
    }


# ==================================================================================================
# FRAMEWORK ORCHESTRATOR
# ==================================================================================================

def run_survival_strategy_framework(
    df: pd.DataFrame,
    config: FrameworkConfig,
    scenarios: Sequence[ScenarioDefinition],
) -> RunArtifacts:
    """Executes the complete governed workflow and returns its immutable artifact package."""
    validation_result = validate_survival_input(df, config, scenarios)
    validation_result.raise_for_errors()
    validation_frame = validation_result.to_frame()

    fingerprint = dataframe_fingerprint(df, config.id_col)
    run_id, run_directory = create_run_directory(config, fingerprint)

    personas = fit_personas(df, config)
    model = fit_cox_model(df, config)
    risk = score_and_stratify_population(df, personas, model, config)
    scenario_artifacts = simulate_scenarios(df, config, scenarios, model, risk)
    model.ph_sensitivity_results = run_ph_sensitivity_analysis(
        df,
        config,
        model,
        risk,
        scenarios_config=scenarios,
        primary_scenario_summary=scenario_artifacts.summary,
    )

    reproducibility_checks = verify_substantive_reproducibility(
        df,
        config,
        scenarios,
        personas,
        model,
        risk,
        scenario_artifacts,
    )

    # Analytical checks are completed before reports so the status and review items can be embedded.
    acceptance_checks = run_acceptance_checks(
        df,
        config,
        personas,
        model,
        risk,
        scenario_artifacts,
        reproducibility_checks,
    )
    acceptance_status = derive_acceptance_status(acceptance_checks)
    narrative = build_executive_narrative(
        config,
        model,
        risk,
        scenario_artifacts,
        acceptance_status,
    )

    plots = generate_evidence_plots(
        run_directory,
        config,
        personas,
        model,
        risk,
        scenario_artifacts,
    )

    power_point = run_directory / "Survival_Strategy_Deck.pptx"
    technical_pdf = run_directory / "Technical_Model_Evidence.pdf"

    export_executive_powerpoint(
        power_point,
        config,
        model,
        personas,
        risk,
        scenario_artifacts,
        plots,
        narrative,
        acceptance_status,
    )
    export_technical_pdf(
        technical_pdf,
        config,
        validation_frame,
        acceptance_checks,
        reproducibility_checks,
        acceptance_status,
        model,
        personas,
        risk,
        scenario_artifacts,
        plots,
        narrative,
    )

    # Add the physical-report check, derive the final status, and regenerate reports so the final
    # status is reflected in both stakeholder artifacts.
    acceptance_checks = run_acceptance_checks(
        df,
        config,
        personas,
        model,
        risk,
        scenario_artifacts,
        reproducibility_checks,
        report_paths=(power_point, technical_pdf),
    )
    acceptance_status = derive_acceptance_status(acceptance_checks)
    narrative = build_executive_narrative(
        config,
        model,
        risk,
        scenario_artifacts,
        acceptance_status,
    )
    export_executive_powerpoint(
        power_point,
        config,
        model,
        personas,
        risk,
        scenario_artifacts,
        plots,
        narrative,
        acceptance_status,
    )
    export_technical_pdf(
        technical_pdf,
        config,
        validation_frame,
        acceptance_checks,
        reproducibility_checks,
        acceptance_status,
        model,
        personas,
        risk,
        scenario_artifacts,
        plots,
        narrative,
    )

    archive_run_tables(
        run_directory,
        config,
        scenarios,
        df,
        validation_frame,
        acceptance_checks,
        reproducibility_checks,
        personas,
        model,
        risk,
        scenario_artifacts,
        narrative,
    )

    metadata = model.fit_metadata.iloc[0]
    run_metadata = {
        "engine_release": ENGINE_RELEASE,
        "run_id": run_id,
        "run_timestamp_utc": utc_now(),
        "project_name": config.project_name,
        "input_fingerprint_sha256": fingerprint,
        "row_count": len(df),
        "event_count": int(df[config.event_col].sum()),
        "event_rate": float(df[config.event_col].mean()),
        "duration_col": config.duration_col,
        "event_col": config.event_col,
        "time_unit_label": config.time_unit_label,
        "id_col": config.id_col,
        "model_features": config.model_features,
        "segmentation_features": config.segmentation_features,
        "risk_quantile": config.risk_quantile,
        "target_cohort_count": len(risk.target_cohort),
        "penalizer": config.penalizer,
        "l1_ratio": config.l1_ratio,
        "random_state": config.random_state,
        "scenario_names": [scenario.name for scenario in scenarios],
        "scenario_types": [scenario.scenario_type for scenario in scenarios],
        "apparent_concordance_index": float(metadata["Apparent_Concordance_Index"]),
        "cv_mean_concordance": float(metadata["CV_Mean_Concordance"]),
        "cv_sd_concordance": float(metadata["CV_SD_Concordance"]),
        "integrated_brier_score": float(metadata["Integrated_Brier_Score"]),
        "ph_global_status": str(metadata["PH_Global_Status"]),
        "persona_silhouette_score": float(personas.quality_metrics.iloc[0]["Silhouette_Score"]),
        "persona_stability_mean_ari": float(personas.quality_metrics.iloc[0]["Stability_Mean_ARI"]),
        "package_versions": package_versions(),
        "acceptance_status": acceptance_status,
        "review_check_count": int(acceptance_checks["Status"].eq("REVIEW").sum()),
        "failed_check_count": int(acceptance_checks["Status"].eq("FAIL").sum()),
    }
    write_json(run_directory / "run_metadata.json", run_metadata)

    manifest = build_manifest(run_directory)
    write_json(run_directory / "artifact_manifest.json", manifest)

    return RunArtifacts(
        run_id=run_id,
        run_directory=run_directory,
        validation=validation_frame,
        personas=personas,
        model=model,
        risk=risk,
        scenarios=scenario_artifacts,
        plots=plots,
        executive_narrative=narrative,
        acceptance_checks=acceptance_checks,
        reproducibility_checks=reproducibility_checks,
        acceptance_status=acceptance_status,
        power_point=power_point,
        technical_pdf=technical_pdf,
        manifest=manifest,
    )

# ==================================================================================================
# ROSSI SMOKE TEST — SECONDARY TECHNICAL DEMONSTRATION
# ==================================================================================================

def build_rossi_smoke_test(
    output_root: Path,
) -> tuple[pd.DataFrame, FrameworkConfig, tuple[ScenarioDefinition, ...]]:
    """Builds a secondary public-data smoke test using the lifelines Rossi dataset.

    Important boundary:
    The two scenarios below are technical sensitivity checks used to prove framework portability.
    They are not presented as retention interventions, causal strategies, or criminal-justice policy.
    """
    data = load_rossi().reset_index(drop=True)
    data.insert(0, "subject_id", [f"ROSSI-{i:04d}" for i in range(1, len(data) + 1)])
    data["age_x_prio"] = data["age"] * data["prio"]

    config = FrameworkConfig(
        project_name="Survival Strategy Framework - Rossi Technical Smoke Test",
        duration_col="week",
        event_col="arrest",
        id_col="subject_id",
        time_unit_label="weeks",
        model_features=(
            "fin",
            "age",
            "race",
            "wexp",
            "mar",
            "paro",
            "prio",
            "age_x_prio",
        ),
        segmentation_features=("age", "prio", "fin"),
        profile_features=("age", "prio", "fin", "wexp", "mar"),
        n_clusters=3,
        risk_quantile=0.75,
        penalizer=0.10,
        random_state=42,
        standardize_model_features=True,
        show_km_confidence_intervals=True,
        run_ph_test=True,
        evaluation_horizon=26.0,
        timeline_points=105,
        output_root=output_root,
        save_input_snapshot=True,
    )

    scenarios = (
        ScenarioDefinition(
            name="No-Change Control - Rossi",
            description=(
                "Apply no changes to verify neutral scenario plumbing in the secondary smoke test."
            ),
            changes=(),
            scenario_type="CONTROL",
            expected_direction="NEUTRAL",
        ),
        ScenarioDefinition(
            name="Financial Support Flag Enabled - Technical Sensitivity",
            description=(
                "Set the public Rossi financial-support indicator to one for the same high-risk "
                "cohort as a technical portability test."
            ),
            changes=(
                FeatureChange(
                    feature="fin",
                    operation="replace",
                    value=1.0,
                    lower_bound=0.0,
                    upper_bound=1.0,
                    round_digits=0,
                ),
            ),
            scenario_type="TECHNICAL",
            expected_direction="UNSPECIFIED",
        ),
        ScenarioDefinition(
            name="Prior Incident Count Lower - Technical Sensitivity",
            description=(
                "Reduce the prior-incident count by one, bounded at zero, solely to verify "
                "dependency-safe same-cohort scoring in a second dataset."
            ),
            changes=(
                FeatureChange(
                    feature="prio",
                    operation="add",
                    value=-1.0,
                    lower_bound=0.0,
                    round_digits=0,
                ),
            ),
            scenario_type="TECHNICAL",
            expected_direction="UNSPECIFIED",
        ),
    )
    return data, config, scenarios


# ==================================================================================================
# BUILT-IN SELF-TESTS
# ==================================================================================================

def run_internal_self_tests() -> pd.DataFrame:
    """Runs lightweight unit-style checks without requiring pytest."""
    results: list[dict[str, str]] = []

    def record(name: str, passed: bool, detail: str) -> None:
        results.append(
            {
                "Test": name,
                "Status": "PASS" if passed else "FAIL",
                "Detail": detail,
            }
        )
        if not passed:
            raise AssertionError(f"Self-test failed: {name} — {detail}")

    data_a = generate_synthetic_retention_data(n_records=1_000, seed=17)
    data_b = generate_synthetic_retention_data(n_records=1_000, seed=17)
    record(
        "Synthetic data reproducibility",
        data_a.equals(data_b),
        "Same seed generated identical data.",
    )

    test_frame = pd.DataFrame(
        {
            "a": [1.0, 2.0],
            "b": [3.0, 4.0],
            "a_x_b": [0.0, 0.0],
            "a_sq": [0.0, 0.0],
        }
    )
    synced, audit = sync_engineered_dependencies(
        test_frame,
        ["a", "b", "a_x_b", "a_sq"],
        "SELF_TEST",
    )
    record(
        "Interaction dependency synchronization",
        np.allclose(synced["a_x_b"], test_frame["a"] * test_frame["b"]),
        "Interaction term was rebuilt from source fields.",
    )
    record(
        "Squared-term dependency synchronization",
        np.allclose(synced["a_sq"], test_frame["a"] ** 2),
        "Squared term was rebuilt from its base field.",
    )
    record(
        "Dependency audit status",
        audit["Validation_Status"].eq("PASS").all(),
        "Post-rebuild audit checks passed.",
    )

    bad_data = data_a.copy()
    bad_data.loc[0, "attrition_event"] = 2
    config = build_default_retention_config(Path("outputs_self_test"))
    validation = validate_survival_input(
        bad_data,
        config,
        build_default_retention_scenarios(),
    )
    record(
        "Nonbinary event rejection",
        validation.has_errors,
        "Invalid event value was rejected before modeling.",
    )

    raw_p = np.array([0.01, 0.03, 0.04])
    holm = holm_adjusted_pvalues(raw_p)
    bh = benjamini_hochberg_pvalues(raw_p)
    record(
        "Multiplicity-adjusted PH p-values",
        np.all(holm >= raw_p) and np.all(bh >= raw_p),
        "Holm and Benjamini-Hochberg adjusted values were not smaller than raw p-values.",
    )
    record(
        "Scenario direction classification",
        (
            modeled_direction(0.0, 0.0, 1e-10) == "NEUTRAL"
            and modeled_direction(0.10, 0.02, 1e-10) == "IMPROVED"
            and modeled_direction(-0.10, -0.02, 1e-10) == "ADVERSE"
        ),
        "Neutral, favorable, and adverse modeled movement were classified correctly.",
    )

    return pd.DataFrame(results)


# ==================================================================================================
# COMMAND-LINE ENTRY POINT
# ==================================================================================================

def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Run the Survival Strategy Framework demonstration.",
    )
    parser.add_argument(
        "--demo",
        choices=("retention", "rossi"),
        default="retention",
        help="Primary synthetic retention demo or secondary Rossi smoke test.",
    )
    parser.add_argument(
        "--records",
        type=int,
        default=7_500,
        help="Synthetic retention row count; ignored for Rossi.",
    )
    parser.add_argument(
        "--seed",
        type=int,
        default=42,
        help="Deterministic synthetic data seed.",
    )
    parser.add_argument(
        "--output-root",
        type=Path,
        default=Path("outputs"),
        help="Root directory for immutable run packages.",
    )
    parser.add_argument(
        "--self-test",
        action="store_true",
        help="Run built-in unit-style checks before the framework run.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()

    try:
        if args.self_test:
            self_test_results = run_internal_self_tests()
            print("\nINTERNAL SELF-TESTS")
            print(self_test_results.to_string(index=False))

        if args.demo == "retention":
            data = generate_synthetic_retention_data(
                n_records=args.records,
                seed=args.seed,
            )
            config = build_default_retention_config(args.output_root)
            scenarios = build_default_retention_scenarios()
        else:
            data, config, scenarios = build_rossi_smoke_test(args.output_root)

        artifacts = run_survival_strategy_framework(
            df=data,
            config=config,
            scenarios=scenarios,
        )

        print("\n" + "=" * 100)
        print("SURVIVAL STRATEGY FRAMEWORK — PIPELINE COMPLETE")
        print("=" * 100)
        print(f"Run ID:             {artifacts.run_id}")
        print(f"Run directory:      {artifacts.run_directory.resolve()}")
        print(f"Executive PPTX:     {artifacts.power_point.resolve()}")
        print(f"Technical PDF:      {artifacts.technical_pdf.resolve()}")
        print(f"Target cohort:      {len(artifacts.risk.target_cohort):,}")
        print("\nScenario results:")
        if artifacts.scenarios.summary.empty:
            print("No configured scenarios.")
        else:
            columns = [
                "Scenario_Rank",
                "Scenario_Name",
                "Scenario_Type",
                "Relative_Hazard_Reduction",
                "Survival_Uplift_Percentage_Points",
                "Modeled_Direction",
            ]
            print(artifacts.scenarios.summary[columns].to_string(index=False))
        print(f"\nRun acceptance:      {artifacts.acceptance_status}")
        print("\nAcceptance checks:")
        print(artifacts.acceptance_checks.to_string(index=False))
        print("\nBoundary: modeled sensitivity, not causal or production impact.")
        return 0

    except Exception as exc:
        print("\nPIPELINE FAILED", file=sys.stderr)
        print(str(exc), file=sys.stderr)
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
